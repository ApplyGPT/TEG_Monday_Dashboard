"""
Proposal Creator
Builds a proposal PPTX from a template-like structure using inputs and checkboxes
"""

import streamlit as st
import os
import sys
from io import BytesIO
import re
from datetime import datetime
from typing import Optional
import base64

# Add project root to path
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
    from pptx.dml.color import RGBColor
except Exception:
    Presentation = None

# Google Slides
try:
    from google.oauth2.service_account import Credentials as SACredentials
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaIoBaseUpload
except Exception:
    SACredentials = None
    build = None
    MediaIoBaseUpload = None

# OAuth imports removed - using service account credentials instead


st.set_page_config(
    page_title="Proposal Creator",
    page_icon="📽️",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Hide tool pages from sidebar
st.markdown(
    """
<style>
    [data-testid="stSidebarNav"] a[href*="tools"] { display: none !important; }
    [data-testid="stSidebarNav"] a[href*="signnow_form"] { display: none !important; }
    [data-testid="stSidebarNav"] a[href*="quickbooks_form"] { display: none !important; }
    [data-testid="stSidebarNav"] a[href*="workbook_creator"] { display: none !important; }
</style>
""",
    unsafe_allow_html=True,
)


def _get_credentials():
    """Load service account credentials from Streamlit secrets (same pattern as google_sheets_uploader)."""
    if not SACredentials:
        raise RuntimeError(
            "Google API libraries not available. Please install google-auth and google-api-python-client."
        )

    info = st.secrets.get("google_service_account")
    if not info:
        raise RuntimeError(
            "Google Cloud service account credentials missing in secrets. "
            "Add `google_service_account` to your Streamlit secrets."
        )

    scopes = [
        "https://www.googleapis.com/auth/drive",
        "https://www.googleapis.com/auth/presentations",
    ]
    return SACredentials.from_service_account_info(
        info,
        scopes=scopes
    )


def create_presentation(title_name: str,
                        priorities_left: list[str],
                        priorities_middle: list[str],
                        priorities_right: list[str],
                        image_bytes: bytes | None) -> bytes:
    """Create a simple PPTX based on inputs. Returns PPTX file bytes.

    Slides:
    1) Title slide with name injected
    2) Priorities slide group A (checkbox set 1-4)
    3) Priorities slide group B (checkbox set 1-3 displayed conditionally)
    4) Image slide (if provided)
    """
    if Presentation is None:
        raise RuntimeError("python-pptx not installed")

    prs = Presentation()

    def add_title_slide(name_text: str):
        slide_layout = prs.slide_layouts[0]  # Title
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Proposal for {name_text}" if name_text else "Proposal"
        subtitle = slide.placeholders[1]
        subtitle.text = f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')}"

    def add_priorities_slide(title_text: str, items: list[str]):
        slide_layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text
        # Add a text box for bullet points
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        for idx, item in enumerate(items):
            if idx == 0:
                p = tf.paragraphs[0]
                p.text = f"• {item}"
            else:
                p = tf.add_paragraph()
                p.text = f"• {item}"
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    def add_image_slide(img_bytes: bytes):
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        # Center image with max width/height
        max_width = Inches(9)
        max_height = Inches(6)
        pic = slide.shapes.add_picture(BytesIO(img_bytes), Inches(1), Inches(1))
        # Scale down if larger than bounds
        if pic.width > max_width or pic.height > max_height:
            ratio_w = max_width / pic.width
            ratio_h = max_height / pic.height
            ratio = min(ratio_w, ratio_h)
            pic.width = int(pic.width * ratio)
            pic.height = int(pic.height * ratio)
            # Re-center
            pic.left = int((prs.slide_width - pic.width) / 2)
            pic.top = int((prs.slide_height - pic.height) / 2)

    # 1) Title
    add_title_slide(title_name)

    # 2) Priorities group A
    if priorities_left:
        add_priorities_slide("Scope Items", priorities_left)

    # 3) Priorities group B (left/middle/right logic simplified per instructions)
    # If only one checked, show it; if multiple, show each on its own slide for now
    groups = [("SOURCING", priorities_middle), ("DEVELOPMENT", priorities_right), ("TREATMENT", priorities_left)]
    selected_groups = [(title, items) for title, items in groups if items]
    if len(selected_groups) == 1:
        add_priorities_slide(selected_groups[0][0], selected_groups[0][1])
    elif len(selected_groups) > 1:
        for title, items in selected_groups:
            add_priorities_slide(title, items)

    # 4) Image slide
    if image_bytes:
        add_image_slide(image_bytes)

    out = BytesIO()
    prs.save(out)
    return out.getvalue()


def create_from_template_pptx(template_path: str,
                              name_value: str,
                              scope_items: list[str],
                              sourcing_items: list[str],
                              development_items: list[str],
                              treatment_items: list[str],
                              image_bytes: bytes | None) -> bytes:
    """Open template PPTX and modify slides to match requested edits.

    Edits required:
    - Slide 2: set name
    - Slide 5: set name and scope items
    - Slide 10: set packages
    - Slide 11: insert image (from PDF first page)
    - Slide 12: set name
    """
    if Presentation is None:
        raise RuntimeError("python-pptx not installed")

    prs = Presentation(template_path)

    def set_title(slide_index: int, text: str, *, force_upper: bool = True, font_name: str | None = "Schibsted Grotesk Medium", font_size_pt: int | None = None):
        if slide_index < 0 or slide_index >= len(prs.slides):
            return
        slide = prs.slides[slide_index]
        # Replace text in-place on the first text_frame without clearing formatting
        def replace_text_preserve_runs(tf, new_text: str):
            if not tf or not tf.paragraphs:
                return
            p = tf.paragraphs[0]
            if not p.runs:
                p.text = new_text
                return
            # Keep formatting of first run
            first_run = p.runs[0]
            # Clear all runs' text
            for r in p.runs:
                r.text = ""
            first_run.text = new_text
            # Apply font
            if font_name:
                first_run.font.name = font_name
            if font_size_pt:
                first_run.font.size = Pt(font_size_pt)
        title_shape = slide.shapes.title
        if title_shape and getattr(title_shape, 'has_text_frame', False):
            replace_text_preserve_runs(title_shape.text_frame, text.upper() if force_upper else text)
            return
        for shape in slide.shapes:
            if getattr(shape, 'has_text_frame', False):
                replace_text_preserve_runs(shape.text_frame, text.upper() if force_upper else text)
                break

    def replace_tokens_on_slide(slide_index: int, replacements: dict[str, str], *, font_name: str | None = "Schibsted Grotesk Medium", font_size_pt: int | None = None, force_upper: bool = True):
        if slide_index < 0 or slide_index >= len(prs.slides):
            return
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            tf = shape.text_frame
            for p in tf.paragraphs:
                if not p.runs:
                    # whole paragraph text fallback
                    t = p.text or ""
                    for old, new in replacements.items():
                        # case-insensitive replacement
                        repl = new.upper() if force_upper else new
                        t = re.sub(re.escape(old), repl, t, flags=re.IGNORECASE)
                    p.text = t
                    # Apply font to paragraph runs
                    for r in p.runs:
                        if font_name:
                            r.font.name = font_name
                        if font_size_pt:
                            r.font.size = Pt(font_size_pt)
                else:
                    for r in p.runs:
                        t = r.text or ""
                        for old, new in replacements.items():
                            repl = new.upper() if force_upper else new
                            t = re.sub(re.escape(old), repl, t, flags=re.IGNORECASE)
                        r.text = t
                        if font_name:
                            r.font.name = font_name
                        if font_size_pt:
                            r.font.size = Pt(font_size_pt)

    def set_body_bullets(slide_index: int, bullets: list[str]):
        if slide_index < 0 or slide_index >= len(prs.slides):
            return
        slide = prs.slides[slide_index]
        # find a likely body placeholder; fallback to add a textbox
        body_shape = None
        for shape in slide.shapes:
            if getattr(shape, 'has_text_frame', False):
                # prefer shapes that are not the title
                if slide.shapes.title and shape == slide.shapes.title:
                    continue
                body_shape = shape
                break
        if body_shape is None:
            body_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = body_shape.text_frame
        tf.clear()
        first = True
        for b in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.text = f"• {b}"
            first = False

    def add_or_replace_image(slide_index: int, image_png_bytes: bytes):
        if slide_index < 0 or slide_index >= len(prs.slides):
            return
        slide = prs.slides[slide_index]
        # Remove existing pictures on the slide
        to_remove = []
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                to_remove.append(shape)
        for shp in to_remove:
            sp = shp._element
            sp.getparent().remove(sp)
        # place centered
        max_width = Inches(9)
        max_height = Inches(6)
        pic = slide.shapes.add_picture(BytesIO(image_png_bytes), Inches(1), Inches(1))
        if pic.width > max_width or pic.height > max_height:
            ratio_w = max_width / pic.width
            ratio_h = max_height / pic.height
            ratio = min(ratio_w, ratio_h)
            pic.width = int(pic.width * ratio)
            pic.height = int(pic.height * ratio)
            pic.left = int((prs.slide_width - pic.width) / 2)
            pic.top = int((prs.slide_height - pic.height) / 2)

    # Utility: normalize text
    def norm(s: str) -> str:
        s = (s or "")
        s = re.sub(r"[^A-Za-z0-9 ]+", " ", s)
        return " ".join(s.strip().split()).upper()

    # Known labels (exact text in template expected)
    scope_labels = [
        "SOURCE FABRIC & TRIMS EFFECTIVELY",
        "DEVELOP HIGH QUALITY PATTERNS & SAMPLES",
        "PRODUCE A SMALL VOLUME PRODUCTION RUN FOR SALES",
        "MANAGE FABRIC TREATMENTS WITH PRECISION",
    ]
    sourcing_labels = [
        "SOURCING",
        "SOURCING INTAKE SESSION",
        "EXPERT INPUT AND PLANNING",
        "SWATCHES AND TRIMS GATHERED",
        "NEGOTIATE PRICING AND MINIMUMS",
        "GUIDANCE IN POs AND ORDERING",
        "TRACKING RECEIPT OF ORDERS",
        "1-2 ROUNDS OF REVISIONS",
    ]
    development_labels = [
        "DEVELOPMENT",
        "DEVELOPMENT ONBOARDING SESSION",
        "TEG SPECIFICATION SHEETS COMPLETED",
        "TECHNICAL INTAKE WITH PATTERN MAKER",
        "FIRST PATTERNS & FIRST SAMPLES",
        "ONE FITTING WITH PATTERN MAKER",
        "ONE ROUND OF FIT ADJUSTMENTS",
        "ONE DUPLICATE PER STYLE",
        "FINAL PRODUCTION READY PATTERNS",
    ]
    treatment_labels = [
        "TREATMENT",
        "TREATMENT INTAKE SESSION",
        "EXPERT INPUT AND PLANNING",
        "ARTWORK / COLOR APPROVAL",
        "NEGOTIATE PRICING AND MINIMUMS",
        "GUIDANCE IN POs AND ORDERING",
        "COORDINATE SEND-OUTS",
        "PROJECT MANAGEMENT",
    ]

    # Slide indices are 0-based
    # Split name
    first_name = (name_value or "").strip().split(" ")[0] if name_value else ""
    last_name = "" if not name_value else " ".join((name_value.strip().split(" ")[1:]))

    # 2: exact full-box replacement only (prevent partial artifacts)
    if name_value:
        slide2 = prs.slides[1] if len(prs.slides) > 1 else None
        if slide2 is not None:
            for shape in slide2.shapes:
                if getattr(shape, 'has_text_frame', False):
                    txt = (shape.text_frame.text or "").strip().upper()
                    if txt == "1ST NAME 2ND NAME":
                        full = f"{first_name} {last_name}".strip().upper()
                        # Clear and rebuild to avoid partial replacement artifacts
                        tf = shape.text_frame
                        tf.clear()
                        p = tf.paragraphs[0]
                        p.text = full
                        if p.runs:
                            p.runs[0].font.name = "Schibsted Grotesk Medium"
                            p.runs[0].font.size = Pt(98.5)
                        break

    # 4: also set name and replace tokens (first name only)
    if name_value:
        replace_tokens_on_slide(3, {"1ST NAME": first_name, "2ND NAME": first_name})

    # 5: set name and toggle scope items by deleting unchecked shapes
    if 4 < len(prs.slides):
        slide5 = prs.slides[4]
        if name_value:
            # replace the specific title token line
            for shape in slide5.shapes:
                if getattr(shape, 'has_text_frame', False):
                    t = (shape.text_frame.text or "").strip().upper()
                    if t in {"1ST NAME’S PRIORITIES", "1ST NAME'S PRIORITIES"}:
                        tf = shape.text_frame
                        tf.clear()
                        p = tf.paragraphs[0]
                        newt = f"{first_name}’S PRIORITIES".upper()
                        if p.runs:
                            p.runs[0].text = newt
                            p.runs[0].font.name = "Schibsted Grotesk Medium"
                        else:
                            p.text = newt
                        break
        desired = set(norm(x) for x in scope_items)
        to_delete = []
        scope_placeholder_shapes = []
        for shape in slide5.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            text = norm(shape.text_frame.text)
            if text in (norm(lbl) for lbl in scope_labels) and text not in desired:
                to_delete.append(shape)
            # Also handle explicit SCOPE ITEM placeholders
            if text in {"SCOPE ITEM 1", "SCOPE ITEM 2", "SCOPE ITEM 3", "SCOPE ITEM 4"}:
                scope_placeholder_shapes.append(shape)
        # Sort placeholders by their position (top-to-bottom) to maintain order
        scope_placeholder_shapes.sort(key=lambda s: s.top)
        
        # Maintain correct order: SOURCE, MANAGE, DEVELOP, PRODUCE
        # But user wants MANAGE as 2nd item, so reorder based on selection
        selected_labels = []
        mapping_list = [
            "SOURCE FABRIC & TRIMS EFFECTIVELY",
            "MANAGE FABRIC TREATMENTS WITH PRECISION",
            "DEVELOP HIGH QUALITY PATTERNS & SAMPLES",
            "PRODUCE A SMALL VOLUME PRODUCTION RUN FOR SALES",
        ]
        # Reorder: MANAGE first if selected, then others in original order
        if norm("MANAGE FABRIC TREATMENTS WITH PRECISION") in desired:
            selected_labels.append("MANAGE FABRIC TREATMENTS WITH PRECISION")
        for lbl in mapping_list:
            if norm(lbl) in desired and lbl != "MANAGE FABRIC TREATMENTS WITH PRECISION":
                selected_labels.append(lbl)
        
        # Assign to placeholders: MANAGE goes to 2nd placeholder (index 1), DEVELOP to 1st (index 0)
        # Reorder selected_labels so MANAGE is in position 1, DEVELOP in position 0
        reordered_labels = []
        for lbl in selected_labels:
            if "MANAGE" in lbl:
                reordered_labels.insert(1, lbl)  # MANAGE goes to position 1
            else:
                reordered_labels.append(lbl)     # Others go in order
        
        for idx, shp in enumerate(scope_placeholder_shapes):
            if idx < len(reordered_labels):
                tf = shp.text_frame
                # replace text, enforce uppercase and font
                if tf.paragraphs:
                    tf.paragraphs[0].text = reordered_labels[idx].upper()
                    run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else None
                    if run:
                        run.font.name = "Schibsted Grotesk Medium"
                        run.font.size = Pt(30)
            else:
                el = shp._element
                el.getparent().remove(el)
        for shp in to_delete:
            el = shp._element
            el.getparent().remove(el)

    # 10: Simple rule - if no package title (font size 32), delete everything below it
    if 9 < len(prs.slides):
        slide10 = prs.slides[9]
        to_delete = []
        
        # First pass: handle package headers and ADD PACKAGE placeholders
        package_headers = {}  # Store position -> package name for headers that exist
        
        for shape in slide10.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            t = norm(shape.text_frame.text)
            
            # Handle existing package headers
            if t in {"SOURCING", "DEVELOPMENT", "TREATMENT"}:
                # Check if this package is selected
                if t == "SOURCING" and sourcing_items:
                    package_headers[shape.left] = "SOURCING"
                elif t == "DEVELOPMENT" and development_items:
                    package_headers[shape.left] = "DEVELOPMENT"
                elif t == "TREATMENT" and treatment_items:
                    package_headers[shape.left] = "TREATMENT"
                else:
                    # This package is not selected, delete the header
                    to_delete.append(shape)
                continue
            
            # Handle ADD PACKAGE placeholders
            if t in {"ADD PACKAGE 1", "ADD PACKAGE 2", "ADD PACKAGE 3"}:
                package_mapping = {
                    "ADD PACKAGE 1": "SOURCING",
                    "ADD PACKAGE 2": "TREATMENT", 
                    "ADD PACKAGE 3": "DEVELOPMENT"
                }
                
                template_package = package_mapping[t]
                
                # Check if this package is selected
                if template_package == "SOURCING" and sourcing_items:
                    shape.text_frame.paragraphs[0].text = "SOURCING"
                    package_headers[shape.left] = "SOURCING"
                elif template_package == "DEVELOPMENT" and development_items:
                    shape.text_frame.paragraphs[0].text = "DEVELOPMENT"
                    package_headers[shape.left] = "DEVELOPMENT"
                elif template_package == "TREATMENT" and treatment_items:
                    shape.text_frame.paragraphs[0].text = "TREATMENT"
                    package_headers[shape.left] = "TREATMENT"
                else:
                    # This package is not selected, delete the placeholder
                    to_delete.append(shape)
                    continue
                
                # enforce font and size
                p = shape.text_frame.paragraphs[0]
                if p.runs:
                    p.runs[0].font.name = "Schibsted Grotesk Medium"
                    p.runs[0].font.size = Pt(32)
                continue
        
        # Second pass: Simple rule - if no header exists for a column, delete everything in that column
        for shape in slide10.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            
            # Skip if already marked for deletion
            if shape in to_delete:
                continue
            
            # Skip headers (already processed)
            t = norm(shape.text_frame.text)
            if t in {"SOURCING", "DEVELOPMENT", "TREATMENT"}:
                continue
            
            # Check if there's a header for this column (based on position)
            shape_left = shape.left
            has_header = False
            
            # Find the closest header to this shape's position
            for header_left in package_headers.keys():
                # If shape is roughly in the same column as a header (within some tolerance)
                if abs(shape_left - header_left) < Inches(2):  # 2 inch tolerance
                    has_header = True
                    break
            
            # If no header exists for this column, delete this shape
            if not has_header:
                to_delete.append(shape)
        
        # Delete all marked shapes
        for shp in to_delete:
            el = shp._element
            el.getparent().remove(el)
    # 11: image
    if image_bytes:
        add_or_replace_image(10, image_bytes)
    # 12: first name possessive (avoid double possessive)
    if first_name:
        replace_tokens_on_slide(11, {"1ST NAME’S": first_name + "’S", "1ST NAME'S": first_name + "’S", "1ST NAME": first_name + "’S", "2ND NAME": ""})

    # Final targeted fixups for pages 2, 5, 10, 12
    # Page 2: replace exact textbox "1ST NAME 2ND NAME"
    if 1 < len(prs.slides) and (first_name or last_name):
        s2 = prs.slides[1]
        for shp in s2.shapes:
            if getattr(shp, 'has_text_frame', False):
                txt = (shp.text_frame.text or "").strip()
                if txt.upper() == "1ST NAME 2ND NAME":
                    new_full = f"{first_name} {last_name}".strip().upper()
                    if shp.text_frame.paragraphs:
                        p = shp.text_frame.paragraphs[0]
                        if p.runs:
                            p.runs[0].text = new_full
                            p.runs[0].font.name = "Schibsted Grotesk Medium"
                        else:
                            p.text = new_full
                    break

    # Page 5: change title to "<FIRST>’S PRIORITIES" and ensure order & font size for placeholders
    if 4 < len(prs.slides) and first_name:
        s5 = prs.slides[4]
        for shp in s5.shapes:
            if getattr(shp, 'has_text_frame', False):
                t = (shp.text_frame.text or "").strip().upper()
                if t in {"1ST NAME’S PRIORITIES", "1ST NAME'S PRIORITIES"}:
                    newt = f"{first_name}’S PRIORITIES".upper()
                    p = shp.text_frame.paragraphs[0]
                    if p.runs:
                        p.runs[0].text = newt
                        p.runs[0].font.name = "Schibsted Grotesk Medium"
                    else:
                        p.text = newt
                    break

        # Re-assert font size 30 on remaining scope placeholders
        for shp in s5.shapes:
            if getattr(shp, 'has_text_frame', False):
                t = (shp.text_frame.text or "").strip()
                if t in {
                    "SOURCE FABRIC & TRIMS EFFECTIVELY",
                    "MANAGE FABRIC TREATMENTS WITH PRECISION",
                    "DEVELOP HIGH QUALITY PATTERNS & SAMPLES",
                    "PRODUCE A SMALL VOLUME PRODUCTION RUN FOR SALES",
                }:
                    p = shp.text_frame.paragraphs[0]
                    if p.runs:
                        p.runs[0].font.name = "Schibsted Grotesk Medium"
                        p.runs[0].font.size = Pt(30)

    # Page 10: if only TREATMENT selected, delete all non-treatment text boxes (incl prices)
    only_treatment = bool(treatment_items) and not (sourcing_items or development_items)
    if only_treatment and 9 < len(prs.slides):
        s10 = prs.slides[9]
        
        # Explicitly define all TREATMENT-related text that should be kept
        treatment_texts = {
            "TREATMENT", "$760 PER SERVICE",
            "TREATMENT INTAKE SESSION", "EXPERT INPUT AND PLANNING",
            "ARTWORK / COLOR APPROVAL", "NEGOTIATE PRICING AND MINIMUMS", 
            "GUIDANCE IN POs AND ORDERING", "COORDINATE SEND-OUTS", "PROJECT MANAGEMENT"
        }
        
        # Also keep partial matches for TREATMENT fields
        treatment_partials = {
            "INTAKE SESSION", "EXPERT INPUT", "ARTWORK", "COLOR APPROVAL",
            "NEGOTIATE PRICING", "GUIDANCE IN", "ORDERING", "COORDINATE",
            "SEND-OUTS", "PROJECT MANAGEMENT"
        }
        
        to_del = []
        headers = []
        
        # Process all text shapes
        for shp in s10.shapes:
            if not getattr(shp, 'has_text_frame', False):
                continue
                
            t = (shp.text_frame.text or "").strip().upper()
            
            # Keep TREATMENT header
            if t == "TREATMENT":
                headers.append(shp)
                continue
                
            # Keep if text exactly matches any TREATMENT text
            if t in treatment_texts:
                continue
                
            # For ambiguous fields that appear in multiple packages, 
            # we need to be more specific about which instances to keep
            # Based on the template analysis, TREATMENT fields are in the rightmost column
            ambiguous_fields = {
                "EXPERT INPUT AND PLANNING", "NEGOTIATE PRICING AND MINIMUMS", 
                "GUIDANCE IN POs AND ORDERING"
            }
            
            if t in ambiguous_fields:
                # Since we're only showing TREATMENT, keep ALL instances of these fields
                # This is safer than trying to guess position
                continue  # Keep this field
                        
            # Keep if text contains any TREATMENT text (for partial matches)
            keep_this = False
            for treatment_text in treatment_texts:
                if treatment_text in t:
                    keep_this = True
                    break
            
            # Also check partial matches
            if not keep_this:
                for partial_text in treatment_partials:
                    if partial_text in t:
                        keep_this = True
                        break
                    
            if not keep_this:
                to_del.append(shp)
        
        # Delete non-TREATMENT text boxes
        for shp in to_del:
            el = shp._element
            el.getparent().remove(el)
            
        # Deduplicate TREATMENT headers - keep only one
        if len(headers) > 1:
            # Keep the first TREATMENT header, delete the rest
            keep_header = headers[0]
            for shp in headers[1:]:
                el = shp._element
                el.getparent().remove(el)

    # Page 12: ensure only FIRST NAME possessive, no double
    if 11 < len(prs.slides) and first_name:
        s12 = prs.slides[11]
        for shp in s12.shapes:
            if getattr(shp, 'has_text_frame', False):
                txt = shp.text_frame.text or ""
                txt2 = re.sub(r"1ST NAME[’']S", f"{first_name}’S", txt, flags=re.IGNORECASE)
                txt2 = re.sub(r"2ND NAME[’']S", f"{first_name}’S", txt2, flags=re.IGNORECASE)
                txt2 = re.sub(r"1ST NAME", f"{first_name}’S", txt2, flags=re.IGNORECASE)
                txt2 = re.sub(r"2ND NAME", "", txt2, flags=re.IGNORECASE)
                txt2 = txt2.replace("'S’S", "’S").replace("’S’S", "’S")
                if txt2 != txt:
                    p = shp.text_frame.paragraphs[0]
                    if p.runs:
                        p.runs[0].text = txt2.upper()
                        p.runs[0].font.name = "Schibsted Grotesk Medium"
                    else:
                        p.text = txt2.upper()

    out = BytesIO()
    prs.save(out)
    return out.getvalue()


def cleanup_old_proposals(drive, parent_folder_id: Optional[str] = None, max_files: int = 10):
    """Clean up old proposal files to prevent quota issues."""
    try:
        # Find all proposal files in the specified folder
        query = "name contains 'Proposal -' and mimeType='application/vnd.google-apps.presentation'"
        if parent_folder_id:
            query += f" and '{parent_folder_id}' in parents"
        
        list_kwargs = {
            "q": query,
            "orderBy": "createdTime desc",
            "fields": "files(id,name,createdTime)",
        }
        
        # Check if folder is in a Shared Drive
        shared_drive_id = None
        if parent_folder_id:
            try:
                folder_info = drive.files().get(
                    fileId=parent_folder_id,
                    fields="driveId",
                    supportsAllDrives=True
                ).execute()
                folder_drive_id = folder_info.get("driveId")
                if folder_drive_id:
                    shared_drive_id = folder_drive_id
                    list_kwargs.update({
                        "supportsAllDrives": True,
                        "includeItemsFromAllDrives": True,
                        "corpora": "drive",
                        "driveId": shared_drive_id,
                    })
            except Exception:
                pass
        
        results = drive.files().list(**list_kwargs).execute()
        files = results.get('files', [])
        
        # If we have more than max_files, delete the oldest ones
        if len(files) > max_files:
            files_to_delete = files[max_files:]
            for file in files_to_delete:
                try:
                    delete_kwargs = {"fileId": file['id']}
                    if shared_drive_id:
                        delete_kwargs["supportsAllDrives"] = True
                    drive.files().delete(**delete_kwargs).execute()
                    st.info(f"Cleaned up old proposal: {file['name']}")
                except Exception as e:
                    st.warning(f"Could not delete {file['name']}: {e}")
    except Exception as e:
        st.warning(f"Could not clean up old files: {e}")


def upload_pptx_to_google_slides(name_value: str,
                                  scope_items: list[str],
                                  sourcing_items: list[str],
                                  development_items: list[str],
                                  treatment_items: list[str],
                                  image_png_bytes: bytes | None) -> str:
    """Create Google Slides from PPTX content using Drive API conversion. Returns new presentationId."""
    if not build or not MediaIoBaseUpload:
        raise RuntimeError("Google API client not available")
    
    # Get credentials using service account (same pattern as google_sheets_uploader)
    creds = _get_credentials()
    drive = build("drive", "v3", credentials=creds)
    
    # Get parent folder ID for proposals from secrets
    cfg = st.secrets.get("google_drive", {}) or {}
    parent_folder_id = cfg.get("parent_folder_id_proposal")
    
    # Clean up old proposal files to prevent quota issues
    cleanup_old_proposals(drive, parent_folder_id)
    
    # Step 1: Generate the perfect PPTX first
    # Get the template path
    inputs_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "inputs")
    template_path = os.path.join(inputs_dir, "Copy of Kim Schultz Dev Deck 10.15.25.pptx")
    
    if os.path.exists(template_path):
        pptx_bytes = create_from_template_pptx(
            template_path, name_value, scope_items, sourcing_items, development_items, treatment_items, image_png_bytes
        )
    else:
        # Fallback to generated presentation if template missing
        pptx_bytes = create_presentation(
            name_value,
            treatment_items,
            sourcing_items,
            development_items,
            image_png_bytes,
        )
    
    # Step 2: Upload PPTX to Google Drive and convert to Google Slides
    # Create media upload from PPTX bytes
    media = MediaIoBaseUpload(
        BytesIO(pptx_bytes), 
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        resumable=False
    )
    
    # File metadata for Google Slides conversion
    presentation_name = f"Proposal - {name_value or 'Client'}"
    file_metadata = {
        "name": presentation_name,
        "mimeType": "application/vnd.google-apps.presentation"  # This tells Drive to convert to Google Slides
    }
    
    # Set parent folder if specified
    if parent_folder_id:
        file_metadata["parents"] = [parent_folder_id]
    
    # Check if folder is in a Shared Drive
    shared_drive_id = None
    if parent_folder_id:
        try:
            folder_info = drive.files().get(
                fileId=parent_folder_id,
                fields="id, name, driveId",
                supportsAllDrives=True
            ).execute()
            folder_drive_id = folder_info.get("driveId")
            if folder_drive_id:
                shared_drive_id = folder_drive_id
        except Exception:
            # If we can't get folder info, continue without Shared Drive support
            pass
    
    try:
        # Upload and convert in one step
        create_kwargs = {
            "body": file_metadata,
            "media_body": media,
            "fields": "id, webViewLink, name"
        }
        # Use supportsAllDrives for Shared Drives
        if shared_drive_id:
            create_kwargs["supportsAllDrives"] = True
        
        uploaded_file = drive.files().create(**create_kwargs).execute()
        
        pres_id = uploaded_file["id"]
        web_link = uploaded_file["webViewLink"]
        
        return pres_id
        
    except Exception as e:
        error_msg = str(e)
        if "quota" in error_msg.lower():
            st.error("❌ Google Drive quota exceeded. Please try again later or check your Google Drive storage.")
        elif "permission" in error_msg.lower():
            st.error("❌ Permission denied. Please check your Google API credentials and permissions.")
        elif "not found" in error_msg.lower():
            st.error("❌ Template file not found. Please ensure the template PPTX file exists in the inputs folder.")
        else:
            st.error(f"❌ Failed to upload and convert PPTX: {error_msg}")
        raise




def main():
    st.title("📽️ Proposal Creator")
    
    inputs_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "inputs")

    # Basic variables
    st.subheader("Variables")
    name_value = st.text_input("Client Name", value=st.session_state.get('pc_name_value', ""), key='pc_name_value')

    # Scope items (exact four items)
    st.subheader("Scope Items")
    s1 = st.checkbox("SOURCE FABRIC & TRIMS EFFECTIVELY", value=st.session_state.get('pc_s1', False), key='pc_s1')
    s2 = st.checkbox("DEVELOP HIGH QUALITY PATTERNS & SAMPLES", value=st.session_state.get('pc_s2', False), key='pc_s2')
    s3 = st.checkbox("PRODUCE A SMALL VOLUME PRODUCTION RUN FOR SALES", value=st.session_state.get('pc_s3', False), key='pc_s3')
    s4 = st.checkbox("MANAGE FABRIC TREATMENTS WITH PRECISION", value=st.session_state.get('pc_s4', False), key='pc_s4')
    group_scope = [label for flag, label in [
        (s1, "SOURCE FABRIC & TRIMS EFFECTIVELY"),
        (s2, "DEVELOP HIGH QUALITY PATTERNS & SAMPLES"),
        (s3, "PRODUCE A SMALL VOLUME PRODUCTION RUN FOR SALES"),
        (s4, "MANAGE FABRIC TREATMENTS WITH PRECISION"),
    ] if flag]

    # Additional packages (three sections)
    st.subheader("Additional Packages")
    colp1, colp2, colp3 = st.columns(3)
    with colp1:
        pkg_sourcing = st.checkbox("SOURCING ($1330 per style)", value=st.session_state.get('pc_pkg_sourcing', False), key='pc_pkg_sourcing')
    with colp2:
        pkg_treatment = st.checkbox("TREATMENT ($760 per service)", value=st.session_state.get('pc_pkg_treatment', False), key='pc_pkg_treatment')
    with colp3:
        pkg_development = st.checkbox("DEVELOPMENT ($2320 per style)", value=st.session_state.get('pc_pkg_development', False), key='pc_pkg_development')

    # Show sub-field descriptions via expanders on page
    with st.expander("SOURCING - sub-fields", expanded=False):
        st.markdown("- SOURCING INTAKE SESSION\n- EXPERT INPUT AND PLANNING\n- SWATCHES AND TRIMS GATHERED\n- NEGOTIATE PRICING AND MINIMUMS\n- GUIDANCE IN POs AND ORDERING\n- TRACKING RECEIPT OF ORDERS\n- 1-2 ROUNDS OF REVISIONS")

    with st.expander("TREATMENT - sub-fields", expanded=False):
        st.markdown("- TREATMENT INTAKE SESSION\n- EXPERT INPUT AND PLANNING\n- ARTWORK / COLOR APPROVAL\n- NEGOTIATE PRICING AND MINIMUMS\n- GUIDANCE IN POs AND ORDERING\n- COORDINATE SEND-OUTS\n- PROJECT MANAGEMENT")

    with st.expander("DEVELOPMENT - sub-fields", expanded=False):
        st.markdown("- DEVELOPMENT ONBOARDING SESSION\n- TEG SPECIFICATION SHEETS COMPLETED\n- TECHNICAL INTAKE WITH PATTERN MAKER\n- FIRST PATTERNS & FIRST SAMPLES\n- ONE FITTING WITH PATTERN MAKER\n- ONE ROUND OF FIT ADJUSTMENTS\n- ONE DUPLICATE PER STYLE\n- FINAL PRODUCTION READY PATTERNS")

    sourcing_items = [
        "SOURCING INTAKE SESSION",
        "EXPERT INPUT AND PLANNING",
        "SWATCHES AND TRIMS GATHERED",
        "NEGOTIATE PRICING AND MINIMUMS",
        "GUIDANCE IN POs AND ORDERING",
        "TRACKING RECEIPT OF ORDERS",
        "1-2 ROUNDS OF REVISIONS",
    ] if pkg_sourcing else []

    development_items = [
        "DEVELOPMENT ONBOARDING SESSION",
        "TEG SPECIFICATION SHEETS COMPLETED",
        "TECHNICAL INTAKE WITH PATTERN MAKER",
        "FIRST PATTERNS & FIRST SAMPLES",
        "ONE FITTING WITH PATTERN MAKER",
        "ONE ROUND OF FIT ADJUSTMENTS",
        "ONE DUPLICATE PER STYLE",
        "FINAL PRODUCTION READY PATTERNS",
    ] if pkg_development else []

    treatment_items = [
        "TREATMENT INTAKE SESSION",
        "EXPERT INPUT AND PLANNING",
        "ARTWORK / COLOR APPROVAL",
        "NEGOTIATE PRICING AND MINIMUMS",
        "GUIDANCE IN POs AND ORDERING",
        "COORDINATE SEND-OUTS",
        "PROJECT MANAGEMENT",
    ] if pkg_treatment else []

    # PDF upload (first page to image)
    st.subheader("PDF Upload")
    pdf_file = st.file_uploader("Upload a PDF to insert as first-page image", type=["pdf"]) 
    img_bytes = st.session_state.get('pc_img_bytes')
    if pdf_file is not None:
        try:
            import fitz  # PyMuPDF
            pdf_bytes = pdf_file.read()
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            if doc.page_count > 0:
                page = doc.load_page(0)
                pix = page.get_pixmap(dpi=180)
                img_bytes = pix.tobytes("png")
                st.session_state['pc_img_bytes'] = img_bytes
                st.success("PDF converted to image for slide")
        except Exception as e:
            st.warning(f"Could not convert PDF to image: {e}")
    elif img_bytes is not None:
        st.info("Using previously uploaded PDF image (restored)")

    # Actions
    col_dl, col_gs = st.columns(2)
    with col_dl:
        generate = st.button("Create PowerPoint", type="primary")
    with col_gs:
        push_gslides = st.button("Create Google Slides", type="primary")
    

    if generate:
        if not Presentation:
            st.error("python-pptx is not installed.")
            return
        try:
            template_path = os.path.join(inputs_dir, "Copy of Kim Schultz Dev Deck 10.15.25.pptx")
            if os.path.exists(template_path):
                pptx_bytes = create_from_template_pptx(
                    template_path,
                    name_value,
                    group_scope,
                    sourcing_items,
                    development_items,
                    treatment_items,
                    img_bytes,
                )
            else:
                # Fallback to generated presentation if template missing
                pptx_bytes = create_presentation(
                    name_value,
                    treatment_items,
                    sourcing_items,
                    development_items,
                    img_bytes,
                )
            
            # Store in session state and trigger auto-download
            st.session_state['pptx_download_data'] = pptx_bytes
            st.session_state['pptx_download_filename'] = f"proposal_{name_value or 'client'}.pptx"
            st.session_state['pptx_download_timestamp'] = datetime.now().timestamp()
            st.session_state['pptx_auto_download'] = True
            st.success("PowerPoint generated and downloaded!!!")
            
        except Exception as e:
            st.error(f"Failed to generate PowerPoint: {e}")
    
    # Auto-download logic: trigger download automatically using data URL
    if st.session_state.get('pptx_auto_download', False):
        pptx_bytes = st.session_state.get('pptx_download_data')
        filename = st.session_state.get('pptx_download_filename', 'proposal.pptx')
        
        if pptx_bytes:
            # Reset the flag immediately to prevent multiple executions
            st.session_state['pptx_auto_download'] = False
            
            # Convert to base64 for data URL download
            b64_pptx = base64.b64encode(pptx_bytes).decode()
            
            # Use Streamlit components for more reliable auto-download
            import streamlit.components.v1 as components
            
            # Use a unique timestamp to ensure component only renders once per generation
            timestamp = st.session_state.get('pptx_download_timestamp', 'default')
            
            components.html(f"""
            <script>
            (function() {{
                // Use a flag with timestamp to ensure download only happens once per generation
                const downloadId = 'pptx_download_{timestamp}';
                if (window.pptxDownloadTriggered === downloadId) {{
                    return;
                }}
                window.pptxDownloadTriggered = downloadId;
                
                function triggerDownload() {{
                    try {{
                        const data = 'data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_pptx}';
                        const link = document.createElement('a');
                        link.href = data;
                        link.download = '{filename}';
                        link.style.display = 'none';
                        document.body.appendChild(link);
                        link.click();
                        setTimeout(() => {{
                            if (link.parentNode) {{
                                document.body.removeChild(link);
                            }}
                        }}, 100);
                        return true;
                    }} catch(e) {{
                        console.error('Download failed:', e);
                        return false;
                    }}
                }}
                
                // Try once with a small delay to ensure DOM is ready
                setTimeout(triggerDownload, 100);
            }})();
            </script>
            """, height=0)
            
            st.success(f"✅ PowerPoint generated: {filename}")

    if push_gslides:
        if not build:
            st.error("Google API client not available")
        else:
            try:
                new_id = upload_pptx_to_google_slides(
                    name_value,
                    group_scope,
                    sourcing_items,
                    development_items,
                    treatment_items,
                    img_bytes,
                )
                st.success(f"Google Slides created successfully! ID: {new_id}")
                st.info(f"View your presentation: https://docs.google.com/presentation/d/{new_id}/edit")
                
            except Exception as e:
                error_msg = str(e)
                if "quota" in error_msg.lower():
                    st.error("❌ Google Drive quota exceeded. Please try again later or check your Google Drive storage.")
                    st.info("💡 **Solutions:**")
                    st.markdown("""
                    - Free up space in your Google Drive
                    - Delete old proposal files from Google Drive
                    - Try using a different Google account with more storage
                    - Wait a few hours and try again (quotas reset periodically)
                    """)
                elif "permission" in error_msg.lower():
                    st.error("❌ Permission denied. Please check your Google API credentials and permissions.")
                elif "not found" in error_msg.lower():
                    st.error("❌ Template file not found. Please ensure the template PPTX file exists in the inputs folder.")
                else:
                    st.error(f"❌ Failed to create Google Slides: {error_msg}")


if __name__ == "__main__":
    main()


