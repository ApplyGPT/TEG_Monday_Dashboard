"""
Deck Creator
Builds a deck PPTX from templates using deck type selection, priorities, and service columns
"""

import streamlit as st
import os
import sys
from io import BytesIO
import re
from datetime import datetime
from typing import Optional, Dict, List, Tuple, Any
import base64
import requests

# Add project root to path
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except Exception:
    Presentation = None
    MSO_ANCHOR = None

# Google Slides
try:
    from google.oauth2.service_account import Credentials as SACredentials
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaIoBaseUpload, MediaIoBaseDownload
except Exception:
    SACredentials = None
    build = None
    MediaIoBaseUpload = None
    MediaIoBaseDownload = None

try:
    from PIL import Image, ImageOps
except Exception:
    Image = None
    ImageOps = None


st.set_page_config(
    page_title="Deck Creator",
    page_icon="📽️",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS to hide all non-tool pages from sidebar navigation
st.markdown(
    """
<style>
/* Hide ALL sidebar list items by default */
[data-testid="stSidebarNav"] li {
    display: none !important;
}

/* Show list items that contain allowed tool pages using :has() selector */
[data-testid="stSidebarNav"] li:has(a[href*="signnow"]),
[data-testid="stSidebarNav"] li:has(a[href*="/tools"]),
[data-testid="stSidebarNav"] li:has(a[href*="workbook"]),
[data-testid="stSidebarNav"] li:has(a[href*="deck_creator"]),
[data-testid="stSidebarNav"] li:has(a[href*="a_la_carte"]) {
    display: block !important;
}

iframe {
    height: 0px !important;
}
</style>
<script>
// JavaScript to show only tool pages and hide everything else
(function() {
    function showToolPagesOnly() {
        const navItems = document.querySelectorAll('[data-testid="stSidebarNav"] li');
        const allowedPages = ['signnow', 'tools', 'workbook', 'deck_creator', 'a_la_carte'];
        
        // Check if we're currently on an ads dashboard page
        const currentUrl = window.location.href.toLowerCase();
        const currentPath = window.location.pathname.toLowerCase();
        const isOnAdsDashboard = currentUrl.includes('ads') && currentUrl.includes('dashboard') ||
                                 currentPath.includes('ads') && currentPath.includes('dashboard');
        
        navItems.forEach(item => {
            const link = item.querySelector('a');
            if (!link) {
                item.style.setProperty('display', 'none', 'important');
                return;
            }
            
            const href = (link.getAttribute('href') || '').toLowerCase();
            const text = link.textContent.trim().toLowerCase();
            
            // Check if this is an allowed tool page
            const isToolPage = allowedPages.some(page => {
                return href.includes(page) || text.includes(page.toLowerCase());
            });
            
            // Make sure it's not ads dashboard or other dashboards
            const isDashboard = (text.includes('ads') && text.includes('dashboard')) || 
                              (text.includes('burki') && text.includes('dashboard')) ||
                              (text.includes('sales') && text.includes('dashboard')) ||
                              (text.includes('database') && text.includes('refresh')) ||
                              (text.includes('new') && text.includes('leads')) ||
                              (text.includes('seo') && text.includes('metrics')) ||
                              (href.includes('ads') && href.includes('dashboard')) ||
                              (href.includes('burki')) ||
                              (href.includes('sales_dashboard')) ||
                              (href.includes('database_refresh')) ||
                              (href.includes('new_leads')) ||
                              (href.includes('seo_metrics'));
            
            // Hide a_la_carte if we're on an ads dashboard page
            const isDevInspection = href.includes('a_la_carte') || text.includes('a_la_carte');
            if (isOnAdsDashboard && isDevInspection) {
                item.style.setProperty('display', 'none', 'important');
                return;
            }
            
            if (isToolPage && !isDashboard) {
                item.style.setProperty('display', 'block', 'important');
                link.style.setProperty('display', 'block', 'important');
            } else {
                item.style.setProperty('display', 'none', 'important');
            }
        });
    }
    
    // Run immediately and on load
    showToolPagesOnly();
    window.addEventListener('load', function() {
        setTimeout(showToolPagesOnly, 50);
        setTimeout(showToolPagesOnly, 200);
        setTimeout(showToolPagesOnly, 500);
    });
    
    // Watch for DOM changes
    const observer = new MutationObserver(function() {
        showToolPagesOnly();
    });
    
    setTimeout(function() {
        const sidebar = document.querySelector('[data-testid="stSidebarNav"]');
        if (sidebar) {
            observer.observe(sidebar, { 
                childList: true, 
                subtree: true,
                attributes: true
            });
        }
    }, 100);
})();
</script>
""",
    unsafe_allow_html=True,
)


# ============================================================================
# Constants and Configuration
# ============================================================================

DECK_TYPES = {
    "Activewear Men": "ACTIVEWEAR MENS DECK.pptx",
    "Activewear Women": "ACTIVEWEAR WOMEN_S DECK.pptx",
    "Bridal": "BRIDAL _  COUTURE DECK_.pptx",
    "Contemporary": "CONTEMPORARY DECK.pptx",
    "Loungewear": "LOUNGEWEAR DECK.pptx",
    "Standards": "STANDARD DECK.pptx",
    "Streetwear": "STREETWEAR DECK.pptx",
}

SERVICE_COLUMNS_FILE = "SERVICE COLUMNS.pptx"
IMAGE_GALLERY_FOLDER_ID = "1Uxn9oWjw22r0F5GPSLIQ_zUTUXfzO5yl"
IMAGE_TARGET_SIZE = (421, 619)  # width x height in px
EMU_PER_INCH = 914400
EMU_PER_PX = EMU_PER_INCH / 96.0
TARGET_WIDTH_EMU = int(IMAGE_TARGET_SIZE[0] * EMU_PER_PX)
TARGET_HEIGHT_EMU = int(IMAGE_TARGET_SIZE[1] * EMU_PER_PX)


# ============================================================================
# Utility Functions
# ============================================================================

def get_inputs_dir() -> str:
    """Get the inputs directory path."""
    return os.path.join(os.path.dirname(os.path.dirname(__file__)), "inputs")


def get_template_path(filename: str) -> str:
    """Get full path to a template file in inputs directory."""
    return os.path.join(get_inputs_dir(), filename)


def normalize_text(text: str) -> str:
    """Normalize text for comparison."""
    if not text:
        return ""
    text = re.sub(r"[^A-Za-z0-9 ]+", " ", text)
    return " ".join(text.strip().split()).upper()


# ============================================================================
# Google API Functions
# ============================================================================

def _get_credentials():
    """Load service account credentials from Streamlit secrets."""
    if not SACredentials:
        raise RuntimeError(
            "Google API libraries not available. Please install google-auth and google-api-python-client."
        )

    info = st.secrets.get("google_service_account")
    if not info:
        raise RuntimeError(
            "Google Cloud service account credentials missing in secrets. "
            "Add `google_service_account` to your Streamlit secrets."
        )

    scopes = [
        "https://www.googleapis.com/auth/drive",
        "https://www.googleapis.com/auth/presentations",
    ]
    return SACredentials.from_service_account_info(
        info,
        scopes=scopes
    )


# ============================================================================
# Google Drive Image Helpers (Gallery)
# ============================================================================

def _get_drive_service():
    """Create a Google Drive service client."""
    if not build:
        raise RuntimeError("Google API client not available")
    creds = _get_credentials()
    return build("drive", "v3", credentials=creds)


@st.cache_data(show_spinner=False)
def list_gallery_images() -> List[Dict[str, str]]:
    """
    List all image files inside the shared gallery folder (including subfolders).
    
    Returns:
        List of dicts: [{"id": str, "name": str, "display_name": str}]
    """
    drive = _get_drive_service()
    drive_cfg = st.secrets.get("google_drive", {}) or {}
    shared_drive_id = drive_cfg.get("shared_drive_id")
    folder_id = IMAGE_GALLERY_FOLDER_ID
    
    # Get root folder name (helps build display path)
    try:
        root_info = drive.files().get(
            fileId=folder_id,
            fields="id, name, driveId",
            supportsAllDrives=True
        ).execute()
        root_name = root_info.get("name", "Gallery")
        root_drive_id = root_info.get("driveId", shared_drive_id)
        if root_drive_id:
            shared_drive_id = root_drive_id
    except Exception:
        root_name = "Gallery"
    
    list_kwargs_base = {
        "supportsAllDrives": True,
        "includeItemsFromAllDrives": True,
        "pageSize": 100,  # Reduced to ensure pagination works correctly
        "fields": "files(id,name,mimeType,parents),nextPageToken",
    }
    if shared_drive_id:
        list_kwargs_base.update({
            "corpora": "drive",
            "driveId": shared_drive_id,
        })
    
    # BFS through folders with pagination support
    queue = [(folder_id, root_name)]
    images: List[Dict[str, str]] = []
    processed_folders = set()  # Track processed folders to avoid infinite loops
    
    def fetch_all_pages(query, list_kwargs):
        """Fetch all pages of results for a query."""
        all_files = []
        page_token = None
        while True:
            kwargs = list_kwargs.copy()
            if page_token:
                kwargs["pageToken"] = page_token
            try:
                response = drive.files().list(q=query, **kwargs).execute()
                page_files = response.get("files", [])
                all_files.extend(page_files)
                page_token = response.get("nextPageToken")
                if not page_token:
                    break
            except Exception as e:
                # Log error but continue
                break
        return all_files
    
    while queue:
        current_id, current_path = queue.pop(0)
        
        # Avoid processing the same folder twice
        if current_id in processed_folders:
            continue
        processed_folders.add(current_id)
        
        try:
            # Subfolders - fetch all pages
            folder_query = f"'{current_id}' in parents and mimeType='application/vnd.google-apps.folder' and trashed=false"
            folder_files = fetch_all_pages(folder_query, list_kwargs_base)
            for f in folder_files:
                folder_id_val = f.get("id")
                folder_name = f.get("name", "Unknown")
                if folder_id_val and folder_id_val not in processed_folders:
                    next_path = f"{current_path}/{folder_name}"
                    queue.append((folder_id_val, next_path))
            
            # Images in current folder - fetch all pages
            image_query = f"'{current_id}' in parents and mimeType contains 'image/' and trashed=false"
            image_files = fetch_all_pages(image_query, list_kwargs_base)
            for img in image_files:
                images.append({
                    "id": img.get("id"),
                    "name": img.get("name") or "image",
                    "display_name": f"{current_path}/{img.get('name') or 'image'}",
                })
        except Exception as e:
            # Log error but continue processing other folders
            continue
    
    # Sort by display name for consistent ordering
    images.sort(key=lambda x: x.get("display_name", "").lower())
    return images


@st.cache_data(show_spinner=False)
def download_gallery_image(file_id: str) -> bytes:
    """Download raw image bytes from Google Drive."""
    if not MediaIoBaseDownload:
        raise RuntimeError("google-api-python-client is required to download images.")
    
    drive = _get_drive_service()
    request = drive.files().get_media(fileId=file_id, supportsAllDrives=True)
    fh = BytesIO()
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    fh.seek(0)
    return fh.read()


def resize_image_to_target(image_bytes: bytes) -> bytes:
    """Resize and letterbox the image to IMAGE_TARGET_SIZE, returning PNG bytes."""
    if not Image or not ImageOps:
        raise RuntimeError("Pillow is required to process images.")
    
    with Image.open(BytesIO(image_bytes)) as img:
        img = ImageOps.exif_transpose(img.convert("RGB"))
        contained = ImageOps.contain(img, IMAGE_TARGET_SIZE, Image.LANCZOS)
        canvas = Image.new("RGB", IMAGE_TARGET_SIZE, (255, 255, 255))
        offset = ((IMAGE_TARGET_SIZE[0] - contained.width) // 2, (IMAGE_TARGET_SIZE[1] - contained.height) // 2)
        canvas.paste(contained, offset)
        out = BytesIO()
        canvas.save(out, format="PNG")
        out.seek(0)
        return out.read()


@st.cache_data(show_spinner=False)
def get_resized_gallery_image(file_id: str) -> bytes:
    """Download and resize a Drive image to the target dimensions."""
    raw = download_gallery_image(file_id)
    return resize_image_to_target(raw)


# ============================================================================
# Slide Extraction Functions (from SERVICE COLUMNS.pptx)
# ============================================================================

@st.cache_data
def extract_slides_from_service_columns() -> Tuple[Dict[int, Dict[str, any]], Dict[int, Dict[str, any]]]:
    """
    Extract priorities and service column slides from SERVICE COLUMNS.pptx.
    
    Returns:
        Tuple of (priorities_dict, service_columns_dict)
        - priorities_dict: {slide_index: {"title": str, "sub_fields": list[str]}}
        - service_columns_dict: {slide_index: {"title": str, "columns": list[str]}}
    """
    if Presentation is None:
        return {}, {}
    
    service_columns_path = get_template_path(SERVICE_COLUMNS_FILE)
    if not os.path.exists(service_columns_path):
        return {}, {}
    
    try:
        prs = Presentation(service_columns_path)
        priorities = {}
        service_columns = {}
        
        for idx, slide in enumerate(prs.slides):
            # Collect all text from the slide
            all_text = []
            all_shapes_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    text = shape.text_frame.text or ""
                    if text.strip():
                        all_text.append(text.strip())
                        all_shapes_text.append(text.strip())
            
            # Normalize all text for analysis
            normalized_text = " ".join([normalize_text(t) for t in all_text])
            
            # Determine slide type based on content
            is_priority = False
            is_service_column = False
            
            priority_keywords = ["PRIORIT", "PRIORITY", "PRIORITIES"]
            service_keywords = ["SERVICE", "COLUMN", "SOURCING", "DEVELOPMENT", "TREATMENT"]
            
            for keyword in priority_keywords:
                if keyword in normalized_text:
                    is_priority = True
                    break
            
            for keyword in service_keywords:
                if keyword in normalized_text:
                    is_service_column = True
                    break
            
            # Extract column names for service columns
            # Look for specific service types: SOURCING, DEVELOPMENT, TREATMENT, DESIGN, WASH/DYE, FABRIC TREATMENT
            if is_service_column or (not is_priority and all_text):
                # Detect services in order of appearance (left to right based on shape position)
                # Collect shapes with their positions
                shape_data = []
                for shape in slide.shapes:
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        text = shape.text_frame.text or ""
                        if text.strip():
                            shape_data.append({
                                'text': text.strip(),
                                'left': shape.left,
                                'top': shape.top
                            })
                
                # Sort by position (left to right, then top to bottom)
                shape_data.sort(key=lambda x: (x['left'], x['top']))
                
                # Service detection patterns
                service_patterns = {
                    "SOURCING": ["SOURCING"],
                    "DEVELOPMENT": ["DEVELOPMENT"],
                    "TREATMENT": ["TREATMENT"],
                    "DESIGN": ["DESIGN"],
                    "WASH/DYE": ["WASH/DYE", "WASH DYE"],
                    "FABRIC TREATMENT": ["FABRIC TREATMENT", "FABRIC"]
                }
                
                found_services = []
                seen_services = set()
                
                # First, collect all text to analyze
                all_text_combined = " ".join([normalize_text(t) for t in all_shapes_text])
                
                # Process shapes in order to detect services (left to right)
                for shape_info in shape_data:
                    text_upper = normalize_text(shape_info['text'])
                    
                    # Check for FABRIC TREATMENT first (more specific)
                    if "FABRIC TREATMENT" in text_upper or ("FABRIC" in text_upper and "TREATMENT" in text_upper):
                        if "FABRIC TREATMENT" not in seen_services:
                            found_services.append("FABRIC TREATMENT")
                            seen_services.add("FABRIC TREATMENT")
                            continue
                    
                    # Check for WASH/DYE
                    if "WASH/DYE" in text_upper or ("WASH" in text_upper and "DYE" in text_upper):
                        if "WASH/DYE" not in seen_services:
                            found_services.append("WASH/DYE")
                            seen_services.add("WASH/DYE")
                            continue
                    
                    # Check for other services in order
                    for service_name, keywords in service_patterns.items():
                        if service_name in seen_services:
                            continue
                        # Skip TREATMENT if FABRIC TREATMENT was already found
                        if service_name == "TREATMENT" and "FABRIC TREATMENT" in seen_services:
                            continue
                        for keyword in keywords:
                            if keyword in text_upper:
                                found_services.append(service_name)
                                seen_services.add(service_name)
                                break
                        if service_name in seen_services:
                            break
                
                # If we didn't find all services from shape order, check the combined text
                # This ensures we catch services even if they're in the same shape or in different order
                all_services_in_text = []
                if "SOURCING" in all_text_combined and "SOURCING" not in seen_services:
                    all_services_in_text.append("SOURCING")
                if "DESIGN" in all_text_combined and "DESIGN" not in seen_services:
                    all_services_in_text.append("DESIGN")
                if ("WASH/DYE" in all_text_combined or ("WASH" in all_text_combined and "DYE" in all_text_combined)) and "WASH/DYE" not in seen_services:
                    all_services_in_text.append("WASH/DYE")
                if ("FABRIC TREATMENT" in all_text_combined or ("FABRIC" in all_text_combined and "TREATMENT" in all_text_combined)) and "FABRIC TREATMENT" not in seen_services:
                    all_services_in_text.append("FABRIC TREATMENT")
                if "TREATMENT" in all_text_combined and "TREATMENT" not in seen_services and "FABRIC TREATMENT" not in seen_services:
                    all_services_in_text.append("TREATMENT")
                if "DEVELOPMENT" in all_text_combined and "DEVELOPMENT" not in seen_services:
                    all_services_in_text.append("DEVELOPMENT")
                
                # If we found services in order, use those. Otherwise, use the combined detection
                if not found_services:
                    found_services = all_services_in_text
                else:
                    # Add any missing services from combined text (maintaining order where possible)
                    for service in all_services_in_text:
                        if service not in found_services:
                            found_services.append(service)
                
                # Create title from found services using dash separator
                if found_services:
                    title = "-".join(found_services)
                else:
                    title = f"Service Columns {idx + 1}"
                
                service_columns[idx] = {
                    "title": title,
                    "columns": found_services
                }
            
            # Extract sub-fields for priorities
            if is_priority:
                # Extract bullet points or list items as sub-fields
                # First, find the title position to only get items below it
                title_bottom = 0
                for shape in slide.shapes:
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        text = normalize_text(shape.text_frame.text or "")
                        if "PRIORIT" in text:
                            title_bottom = shape.top + shape.height
                            break
                
                # Collect shapes with their positions to maintain order
                priority_shapes = []
                for shape in slide.shapes:
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        text = shape.text_frame.text or ""
                        text_clean = text.strip()
                        text_normalized = normalize_text(text)
                        
                        # Skip the title itself
                        if "PRIORIT" in text_normalized:
                            continue
                        
                        # Only include shapes that are below the title (or all if no title found)
                        if text_clean and (title_bottom == 0 or shape.top > title_bottom):
                            priority_shapes.append({
                                'text': text_clean,
                                'left': shape.left,
                                'top': shape.top
                            })
                
                # Sort by position (top to bottom, then left to right)
                priority_shapes.sort(key=lambda x: (x['top'], x['left']))
                
                sub_fields = []
                # Common patterns to exclude (headers, titles, etc.)
                exclude_patterns = [
                    "PRIORIT", "PRIORITY", "PRIORITIES", 
                    "SERVICE", "COLUMN"
                ]
                
                for shape_info in priority_shapes:
                    text = shape_info['text']
                    text_normalized = normalize_text(text)
                    
                    # Skip if it matches any exclude pattern
                    should_exclude = False
                    for pattern in exclude_patterns:
                        if pattern in text_normalized:
                            should_exclude = True
                            break
                    
                    if not should_exclude:
                        # Clean up the text and add as sub-field
                        cleaned = text.strip()
                        # Remove bullet points if present
                        cleaned = re.sub(r'^[•\-\*]\s*', '', cleaned)
                        # Remove extra whitespace
                        cleaned = re.sub(r'\s+', ' ', cleaned)
                        # Skip very short text, empty text, or single characters
                        if cleaned and len(cleaned) > 5:  # Increased minimum length to 5
                            # Additional check: skip if it's just a single word that's likely a header
                            words = cleaned.split()
                            if len(words) > 1 or (len(words) == 1 and len(cleaned) > 10):
                                sub_fields.append(cleaned)
                
                # Number priorities as "Priorities 1", "Priorities 2", etc.
                priority_num = len(priorities) + 1
                priorities[idx] = {
                    "title": f"Priorities {priority_num}",
                    "sub_fields": sub_fields
                }
        
        return priorities, service_columns
    except Exception as e:
        st.warning(f"Could not extract slides from SERVICE COLUMNS.pptx: {e}")
        return {}, {}


# ============================================================================
# Slide Manipulation Functions
# ============================================================================

def replace_slide_with_source(target_prs: Presentation, target_slide_index: int, source_prs: Presentation, source_slide_index: int):
    """
    Replace a slide in target presentation by clearing all shapes and copying from source.
    This is safer than deleting/inserting as it preserves the slide structure.
    
    Args:
        target_prs: Target presentation
        target_slide_index: Index of slide to replace (0-based)
        source_prs: Source presentation
        source_slide_index: Index of slide to copy from source (0-based)
    """
    if target_slide_index < 0 or target_slide_index >= len(target_prs.slides):
        raise IndexError(f"Target slide index {target_slide_index} out of range")
    if source_slide_index < 0 or source_slide_index >= len(source_prs.slides):
        raise IndexError(f"Source slide index {source_slide_index} out of range")
    
    # Get the target slide (don't delete it - preserve structure)
    target_slide = target_prs.slides[target_slide_index]
    source_slide = source_prs.slides[source_slide_index]
    
    # Remove all existing shapes from target slide
    # Get all shapes first, then remove them
    shapes_to_remove = []
    for shape in target_slide.shapes:
        shapes_to_remove.append(shape)
    
    # Remove shapes from the XML
    for shape in shapes_to_remove:
        try:
            sp = shape._element
            parent = sp.getparent()
            if parent is not None:
                parent.remove(sp)
        except (AttributeError, Exception):
            # If removal fails, try alternative method
            try:
                # Try removing via shape's parent
                if hasattr(shape, '_element'):
                    el = shape._element
                    el.getparent().remove(el)
            except Exception:
                # If all else fails, continue
                pass
    
    # Now copy all shapes from source slide to the target slide
    # Process shapes in order to maintain layout
    for shape in source_slide.shapes:
        try:
            # Skip placeholders from slide layout to avoid duplication
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                continue
            
            # Get shape properties
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
                
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                # Copy picture
                try:
                    image_stream = BytesIO()
                    image = shape.image
                    image_stream.write(image.blob)
                    image_stream.seek(0)
                    target_slide.shapes.add_picture(
                        image_stream,
                        left,
                        top,
                        width,
                        height
                    )
                except Exception:
                    # Silently skip if we can't copy picture
                    continue
            elif shape.shape_type == 1:  # MSO_SHAPE_TYPE.LINE
                # Copy line shape using XML for maximum fidelity
                try:
                    source_element = shape._element
                    target_spTree = target_slide.shapes._spTree
                    from lxml import etree
                    new_element = etree.fromstring(etree.tostring(source_element))
                    target_spTree.append(new_element)
                except Exception:
                    # Silently skip if we can't copy line
                    continue
            elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                # Copy text box with formatting
                try:
                    textbox = target_slide.shapes.add_textbox(left, top, width, height)
                    target_tf = textbox.text_frame
                    source_tf = shape.text_frame
                    
                    # Copy shape fill properties (background color, etc.)
                    try:
                        if hasattr(shape, 'fill') and shape.fill.type:
                            if shape.fill.type == 1:  # Solid fill
                                if hasattr(shape.fill, 'fore_color') and hasattr(shape.fill.fore_color, 'rgb'):
                                    textbox.fill.solid()
                                    textbox.fill.fore_color.rgb = shape.fill.fore_color.rgb
                    except Exception:
                        pass
                    
                    # Copy shape line properties (borders, etc.)
                    try:
                        if hasattr(shape, 'line') and hasattr(shape.line, 'color'):
                            if hasattr(shape.line.color, 'rgb') and shape.line.color.rgb:
                                textbox.line.color.rgb = shape.line.color.rgb
                        if hasattr(shape, 'line') and hasattr(shape.line, 'width'):
                            textbox.line.width = shape.line.width
                    except Exception:
                        pass
                    
                    # Enable word wrap for text frames (important for slide 5)
                    target_tf.word_wrap = True
                    
                    # For slide 10 (Services), set vertical alignment to middle
                    if target_slide_index == 9:  # Slide 10 (0-based index 9)
                        try:
                            if MSO_ANCHOR:
                                target_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                        except Exception:
                            pass
                    
                    # Clear default paragraph
                    target_tf.clear()
                    
                    # Copy paragraphs with all formatting
                    for src_para in source_tf.paragraphs:
                        tgt_para = target_tf.add_paragraph()
                        # Set text first, then formatting
                        para_text = src_para.text or ""
                        
                        # Copy paragraph properties
                        tgt_para.alignment = src_para.alignment
                        tgt_para.level = src_para.level
                        if hasattr(src_para, 'space_after'):
                            tgt_para.space_after = src_para.space_after
                        if hasattr(src_para, 'space_before'):
                            tgt_para.space_before = src_para.space_before
                        
                        # Copy runs with formatting (if runs exist)
                        if src_para.runs:
                            # Build text from runs to preserve formatting
                            for src_run in src_para.runs:
                                tgt_run = tgt_para.add_run()
                                tgt_run.text = src_run.text or ""
                                
                                # Copy font properties
                                if src_run.font.name:
                                    tgt_run.font.name = src_run.font.name
                                if src_run.font.size:
                                    tgt_run.font.size = src_run.font.size
                                tgt_run.font.bold = src_run.font.bold
                                tgt_run.font.italic = src_run.font.italic
                                tgt_run.font.underline = src_run.font.underline
                                
                                # Copy color safely
                                if src_run.font.color:
                                    try:
                                        if hasattr(src_run.font.color, 'rgb') and src_run.font.color.rgb:
                                            tgt_run.font.color.rgb = src_run.font.color.rgb
                                    except (AttributeError, TypeError):
                                        pass
                        else:
                            # No runs, just set text
                            tgt_para.text = para_text
                except Exception as e:
                    # Log error but continue - some text boxes may not be copyable
                    continue
            else:
                # For other shape types (lines, connectors, etc.), copy using XML for maximum fidelity
                try:
                    source_element = shape._element
                    target_spTree = target_slide.shapes._spTree
                    from lxml import etree
                    new_element = etree.fromstring(etree.tostring(source_element))
                    target_spTree.append(new_element)
                except Exception:
                    # Silently skip if we can't copy
                    pass
        except Exception:
            # Silently continue - some shapes may not be copyable
            continue


def copy_slide_from_source(target_prs: Presentation, source_prs: Presentation, source_slide_index: int) -> int:
    """
    Copy a slide from source presentation to target presentation.
    Uses XML copying for better fidelity.
    
    Returns:
        Index of the newly added slide in target_prs
    """
    source_slide = source_prs.slides[source_slide_index]
    
    # Use blank layout as base - check if layout 6 exists, otherwise use layout 0
    try:
        blank_layout = target_prs.slide_layouts[6]  # Blank layout
    except IndexError:
        # If layout 6 doesn't exist, try to find a blank layout or use the first one
        blank_layout = None
        for layout in target_prs.slide_layouts:
            if "blank" in layout.name.lower() or layout.name == "Blank":
                blank_layout = layout
                break
        if blank_layout is None:
            blank_layout = target_prs.slide_layouts[0]  # Fallback to first layout
    
    new_slide = target_prs.slides.add_slide(blank_layout)
    
    # Copy all shapes from source to target
    for shape in source_slide.shapes:
        try:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                # Copy picture
                image_stream = BytesIO()
                image = shape.image
                image_stream.write(image.blob)
                image_stream.seek(0)
                new_slide.shapes.add_picture(
                    image_stream,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )
            elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                # Copy text box with formatting
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                textbox = new_slide.shapes.add_textbox(left, top, width, height)
                target_tf = textbox.text_frame
                source_tf = shape.text_frame
                
                # Clear default paragraph
                target_tf.clear()
                
                # Copy paragraphs
                for src_para in source_tf.paragraphs:
                    tgt_para = target_tf.add_paragraph()
                    tgt_para.text = src_para.text
                    tgt_para.alignment = src_para.alignment
                    tgt_para.level = src_para.level
                    
                    # Copy runs with formatting
                    for src_run in src_para.runs:
                        tgt_run = tgt_para.add_run()
                        tgt_run.text = src_run.text
                        if src_run.font.name:
                            tgt_run.font.name = src_run.font.name
                        if src_run.font.size:
                            tgt_run.font.size = src_run.font.size
                        tgt_run.font.bold = src_run.font.bold
                        tgt_run.font.italic = src_run.font.italic
                        tgt_run.font.underline = src_run.font.underline
                        # Fix color error - check if color exists and has rgb property
                        if src_run.font.color:
                            try:
                                if hasattr(src_run.font.color, 'rgb') and src_run.font.color.rgb:
                                    tgt_run.font.color.rgb = src_run.font.color.rgb
                            except (AttributeError, TypeError):
                                # Color doesn't have rgb property, skip it
                                pass
            else:
                # For other shape types, try to copy as picture or skip
                # This is a simplified approach - complex shapes may not copy perfectly
                pass
        except Exception as e:
            # Log but continue - some shapes may not be copyable
            st.warning(f"Could not copy shape: {e}")
            continue
    
    return len(target_prs.slides) - 1


def replace_text_in_slide(slide, old_text: str, new_text: str, case_sensitive: bool = False):
    """Replace text in all text frames of a slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            tf = shape.text_frame
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    if case_sensitive:
                        if old_text in text:
                            run.text = text.replace(old_text, new_text)
                    else:
                        if old_text.lower() in text.lower():
                            # Preserve case of original text
                            pattern = re.compile(re.escape(old_text), re.IGNORECASE)
                            run.text = pattern.sub(new_text, text)


def set_proposal_for_name(slide, client_name: str):
    """Set the client name in the 'Proposal For' slide (slide 2).
    Only replaces '1ST NAME 2ND NAME', keeps 'PROPOSAL FOR' unchanged.
    """
    if not client_name:
        return
    
    first_name = client_name.strip().split(" ")[0] if client_name else ""
    last_name = " ".join(client_name.strip().split(" ")[1:]) if client_name else ""
    full_name = f"{first_name} {last_name}".strip()
    
    # Only replace "1ST NAME 2ND NAME", not "PROPOSAL FOR"
    for shape in slide.shapes:
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            txt = (shape.text_frame.text or "").strip().upper()
            if txt == "1ST NAME 2ND NAME":
                # Clear and rebuild to avoid partial replacement artifacts
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = full_name.upper()
                if p.runs:
                    p.runs[0].font.name = "Schibsted Grotesk Medium"
                    p.runs[0].font.size = Pt(86)
                break


def add_or_replace_image(prs: Presentation, slide_index: int, image_png_bytes: bytes):
    """Add or replace image on a slide, filling 100% of the slide space."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        return
    slide = prs.slides[slide_index]
    # Remove existing pictures on the slide
    to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            to_remove.append(shape)
    for shp in to_remove:
        sp = shp._element
        sp.getparent().remove(sp)
    
    # Add picture and resize to fill 100% of slide (width and height)
    pic = slide.shapes.add_picture(BytesIO(image_png_bytes), 0, 0)
    
    # Set image to fill entire slide
    pic.left = 0
    pic.top = 0
    pic.width = prs.slide_width
    pic.height = prs.slide_height


def populate_gallery_slide(prs: Presentation, images: List[Dict[str, Any]], slide_index: int = 8):
    """
    Populate slide 9 (index 8) with three gallery images.
    
    Args:
        prs: Presentation object
        images: List of dicts [{"id": str, "name": str, "data": bytes}]
        slide_index: Index of the slide to populate (default 8)
    """
    if not images or len(images) != 3:
        return
    if slide_index < 0 or slide_index >= len(prs.slides):
        st.warning("Template doesn't have a 9th slide to place gallery images.")
        return
    
    slide = prs.slides[slide_index]
    
    # Find "PROJECT HIGHLIGHT" text to position images below it
    project_highlight_bottom = None
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = (shape.text_frame.text or "").strip().upper()
            if "PROJECT HIGHLIGHT" in text:
                # Get the bottom position of the PROJECT HIGHLIGHT text
                project_highlight_bottom = shape.top + shape.height
                break
    
    # Find placeholders to replace their text (but we'll use calculated positions)
    placeholder_map = {}
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = (shape.text_frame.text or "").strip().upper()
            if text in {"PLACEHOLDER 1", "PLACEHOLDER 2", "PLACEHOLDER 3"}:
                placeholder_map[text] = shape
    
    # Calculate positions below PROJECT HIGHLIGHT, centered horizontally
    if project_highlight_bottom is not None:
        # Position images below PROJECT HIGHLIGHT with some spacing
        spacing_below_title = int(0.3 * EMU_PER_INCH)  # 0.3 inches below title
        image_top = project_highlight_bottom + spacing_below_title
        
        # Calculate horizontal positions (centered, with gaps between images)
        total_images_width = 3 * TARGET_WIDTH_EMU
        gap = int((prs.slide_width - total_images_width) / 4) if prs.slide_width > total_images_width else int(0.2 * EMU_PER_INCH)
        gap = max(gap, int(0.2 * EMU_PER_INCH))  # Minimum gap of 0.2 inches
        
        left_start = gap
        positions = []
        for idx in range(3):
            left = left_start + idx * (TARGET_WIDTH_EMU + gap)
            # Get placeholder shape if it exists for text replacement
            key = f"PLACEHOLDER {idx + 1}"
            placeholder_shape = placeholder_map.get(key)
            positions.append((left, image_top, TARGET_WIDTH_EMU, TARGET_HEIGHT_EMU, placeholder_shape))
    else:
        # Fallback: center images vertically if PROJECT HIGHLIGHT not found
        gap = int((prs.slide_width - 3 * TARGET_WIDTH_EMU) / 4) if prs.slide_width > 0 else 0
        gap = max(gap, 0)
        top = int((prs.slide_height - TARGET_HEIGHT_EMU) / 2) if prs.slide_height else 0
        positions = []
        left = gap
        for idx in range(3):
            key = f"PLACEHOLDER {idx + 1}"
            placeholder_shape = placeholder_map.get(key)
            positions.append((left, top, TARGET_WIDTH_EMU, TARGET_HEIGHT_EMU, placeholder_shape))
            left += TARGET_WIDTH_EMU + gap
    
    # Place each image and caption
    for idx, img in enumerate(images):
        left, top, width, height, placeholder_shape = positions[idx]
        
        if placeholder_shape:
            try:
                # Remove trailing numbers from name for placeholder text
                display_name = _remove_trailing_numbers(img.get("name") or "")
                placeholder_shape.text_frame.text = display_name.upper()
            except Exception:
                pass
            try:
                sp = placeholder_shape._element
                sp.getparent().remove(sp)
            except Exception:
                pass
        
        # Add picture and set dimensions explicitly to prevent compression/distortion
        # Images are pre-resized to exactly 421x619px, so we set both dimensions explicitly
        pic = slide.shapes.add_picture(BytesIO(img.get("data")), left, top)
        
        # Always use target dimensions to ensure images aren't compressed
        # Setting both width and height is safe because images are exactly 421x619px
        pic.width = TARGET_WIDTH_EMU
        pic.height = TARGET_HEIGHT_EMU
        
        # Add caption with image name below the picture
        caption_top = top + pic.height + int(0.08 * EMU_PER_INCH)
        caption_height = int(0.4 * EMU_PER_INCH)
        try:
            cap_box = slide.shapes.add_textbox(left, caption_top, pic.width, caption_height)
            cap_tf = cap_box.text_frame
            cap_tf.clear()
            cap_p = cap_tf.paragraphs[0]
            # Remove trailing numbers from name for caption
            display_name = _remove_trailing_numbers(img.get("name") or "")
            cap_p.text = display_name.strip()
            cap_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            if cap_p.runs:
                cap_p.runs[0].font.name = "Schibsted Grotesk Medium"
                cap_p.runs[0].font.size = Pt(18)
        except Exception:
            # If caption fails, just continue without breaking deck creation
            pass


# ============================================================================
# Deck Creation Functions
# ============================================================================

def create_deck_from_template(
    deck_type: str,
    client_name: str,
    priority_slide_index: Optional[int],
    service_column_slide_index: Optional[int],
    image_bytes: bytes | None = None,
    gallery_images: Optional[List[Dict[str, bytes]]] = None
) -> bytes:
    """
    Create a deck from the selected template with customizations.
    
    Args:
        deck_type: Selected deck type (key from DECK_TYPES)
        client_name: Client name for "Proposal For" slide
        priority_slide_index: Index of priority slide to insert (from SERVICE COLUMNS.pptx)
        service_column_slide_index: Index of service column slide to insert (from SERVICE COLUMNS.pptx)
        image_bytes: Optional PDF first-page image bytes to place on slide 11
        gallery_images: Optional list of three resized gallery images for slide 9
    
    Returns:
        PPTX file bytes
    """
    if Presentation is None:
        raise RuntimeError("python-pptx not installed")
    
    # Get template path
    template_filename = DECK_TYPES.get(deck_type)
    if not template_filename:
        raise ValueError(f"Unknown deck type: {deck_type}")
    
    template_path = get_template_path(template_filename)
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")
    
    # Open template
    prs = Presentation(template_path)
    
    # Step 1: Set "Proposal For" slide (slide 2, index 1)
    # Replace "1ST NAME 2ND NAME" with client name, keep "PROPOSAL FOR" unchanged
    if len(prs.slides) > 1 and client_name:
        slide2 = prs.slides[1]
        set_proposal_for_name(slide2, client_name)
    
    # Step 1b: Set slide 4 (index 3) - replace "1ST NAME'S PRIORITIES" with first name only
    if len(prs.slides) > 3 and client_name:
        first_name = client_name.strip().split(" ")[0] if client_name else ""
        slide4 = prs.slides[3]
        # First try exact match replacement
        for shape in slide4.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                t = (shape.text_frame.text or "").strip().upper()
                # Check for both apostrophe styles
                if t in {"1ST NAME'S PRIORITIES", "1ST NAME'S PRIORITIES"}:
                    tf = shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    newt = f"{first_name}'S PRIORITIES".upper()
                    if p.runs:
                        p.runs[0].text = newt
                        p.runs[0].font.name = "Schibsted Grotesk Medium"
                    else:
                        p.text = newt
                    break
        
        # Also use token replacement as fallback
        # This handles cases where the text might be in a different format
        # Fix double possessive issue like "NAME'S'S" -> "NAME'S"
        for shape in slide4.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            tf = shape.text_frame
            for p in tf.paragraphs:
                if not p.runs:
                    # whole paragraph text fallback
                    t = p.text or ""
                    # Replace tokens
                    t = re.sub(r"1ST NAME[’']?S", f"{first_name}'S", t, flags=re.IGNORECASE)
                    t = re.sub(r"1ST NAME", f"{first_name}'S", t, flags=re.IGNORECASE)
                    # Fix double possessive
                    t = t.replace("'S'S", "'S").replace("'S'S", "'S")
                    p.text = t.upper()
                    # Apply font to paragraph runs
                    for r in p.runs:
                        r.font.name = "Schibsted Grotesk Medium"
                else:
                    for r in p.runs:
                        t = r.text or ""
                        # Replace tokens
                        t = re.sub(r"1ST NAME[’']?S", f"{first_name}'S", t, flags=re.IGNORECASE)
                        t = re.sub(r"1ST NAME", f"{first_name}'S", t, flags=re.IGNORECASE)
                        # Fix double possessive
                        t = t.replace("'S'S", "'S").replace("'S'S", "'S")
                        r.text = t.upper()
                        r.font.name = "Schibsted Grotesk Medium"
    
    # Step 2: Replace slide 5 (index 4) with priority slide if selected
    if priority_slide_index is not None:
        service_columns_path = get_template_path(SERVICE_COLUMNS_FILE)
        if os.path.exists(service_columns_path):
            try:
                source_prs = Presentation(service_columns_path)
                if 0 <= priority_slide_index < len(source_prs.slides):
                    if len(prs.slides) > 4:
                        # Replace slide 5 (index 4)
                        replace_slide_with_source(prs, 4, source_prs, priority_slide_index)
                        
                        # Add "PRIORITIES" title to slide 5 if it doesn't exist and format it
                        slide5 = prs.slides[4]
                        has_priorities_title = False
                        title_shape = None
                        
                        # First, find or create the title
                        for shape in slide5.shapes:
                            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                                text = (shape.text_frame.text or "").strip().upper()
                                if text == "PRIORITIES":
                                    has_priorities_title = True
                                    title_shape = shape
                                    break
                        
                        if not has_priorities_title:
                            # Center the title horizontally on the slide
                            title_width = Inches(8)
                            title_left = int((prs.slide_width - title_width) / 2)
                            title_box = slide5.shapes.add_textbox(title_left, Inches(0.5), title_width, Inches(0.8))
                            title_tf = title_box.text_frame
                            title_tf.clear()
                            title_p = title_tf.paragraphs[0]
                            title_p.text = "PRIORITIES"
                            title_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                            if title_p.runs:
                                title_p.runs[0].font.name = "Schibsted Grotesk Medium"
                                title_p.runs[0].font.size = Pt(63.5)
                                title_p.runs[0].font.bold = False
                            else:
                                # If no runs, create one
                                run = title_p.add_run()
                                run.text = "PRIORITIES"
                                run.font.name = "Schibsted Grotesk Medium"
                                run.font.size = Pt(63.5)
                                run.font.bold = False
                        else:
                            # Update existing title formatting and center it
                            if title_shape and title_shape.has_text_frame:
                                # Center the shape horizontally
                                title_width = title_shape.width
                                title_shape.left = int((prs.slide_width - title_width) / 2)
                                
                                # Update text formatting
                                for para in title_shape.text_frame.paragraphs:
                                    para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                                    if para.runs:
                                        for run in para.runs:
                                            run.font.name = "Schibsted Grotesk Medium"
                                            run.font.size = Pt(63.5)
                                            run.font.bold = False
                                    else:
                                        # If no runs, update paragraph directly
                                        para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                                        run = para.add_run()
                                        run.text = para.text
                                        run.font.name = "Schibsted Grotesk Medium"
                                        run.font.size = Pt(63.5)
                                        run.font.bold = False
                        
                        # Ensure word wrap is enabled on all text frames in slide 5
                        # Find the PRIORITIES title to determine what's below it
                        priorities_title_shape = None
                        for shape in slide5.shapes:
                            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                                text = (shape.text_frame.text or "").strip().upper()
                                if text == "PRIORITIES":
                                    priorities_title_shape = shape
                                    break
                        
                        # Get the bottom position of the title (top + height)
                        title_bottom = 0
                        if priorities_title_shape:
                            title_bottom = priorities_title_shape.top + priorities_title_shape.height
                        
                        # Fill all boxes below PRIORITIES with background color #c9a47a (RGB: 201, 164, 122)
                        fill_color = RGBColor(201, 164, 122)
                        for shape in slide5.shapes:
                            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                                text = (shape.text_frame.text or "").strip().upper()
                                
                                # Skip the title itself
                                if text == "PRIORITIES":
                                    continue
                                
                                # Enable word wrap for all text frames
                                shape.text_frame.word_wrap = True
                                
                                # If this box is below the title, apply background color
                                if shape.top > title_bottom:
                                    try:
                                        shape.fill.solid()
                                        shape.fill.fore_color.rgb = fill_color
                                    except Exception:
                                        # If setting fill fails, try alternative method
                                        try:
                                            fill = shape.fill
                                            fill.solid()
                                            fill.fore_color.rgb = fill_color
                                        except Exception:
                                            pass
                    else:
                        st.warning("Template doesn't have a 5th slide to replace")
            except Exception as e:
                st.warning(f"Could not replace priority slide: {e}")
    
    # Step 3: Populate slide 9 (index 8) with gallery images
    if gallery_images:
        try:
            populate_gallery_slide(prs, gallery_images, slide_index=8)
        except Exception as e:
            st.warning(f"Could not populate gallery images on slide 9: {e}")
    
    # Step 4: Copy entire slide content from SERVICE COLUMNS.pptx to slide 10 (index 9)
    # Since slide 10 is now empty, we just copy the entire selected slide
    if service_column_slide_index is not None:
        service_columns_path = get_template_path(SERVICE_COLUMNS_FILE)
        if os.path.exists(service_columns_path):
            try:
                source_prs = Presentation(service_columns_path)
                if 0 <= service_column_slide_index < len(source_prs.slides):
                    if len(prs.slides) > 9:
                        # Copy entire slide content from SERVICE COLUMNS.pptx to slide 10
                        replace_slide_with_source(prs, 9, source_prs, service_column_slide_index)
                        
                        # Ensure vertical middle alignment for all text frames in slide 10
                        slide10 = prs.slides[9]
                        for shape in slide10.shapes:
                            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                                try:
                                    if MSO_ANCHOR:
                                        shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                                except Exception:
                                    pass
                    else:
                        st.warning("Template doesn't have a 10th slide")
            except Exception as e:
                st.warning(f"Could not copy service column slide: {e}")
    
    # Step 5: Set slide 12 (index 11) - replace "1ST NAME'S APPROVAL OF PROJECT" with first name only
    if len(prs.slides) > 11 and client_name:
        first_name = client_name.strip().split(" ")[0] if client_name else ""
        slide12 = prs.slides[11]
        # Use regex replacement to handle various formats, similar to slide 4
        for shape in slide12.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            tf = shape.text_frame
            for p in tf.paragraphs:
                if p.runs:
                    for r in p.runs:
                        t = r.text or ""
                        # Replace "1ST NAME'S" or "1ST NAME" with first name + 'S
                        t = re.sub(r"1ST NAME[’']?S", f"{first_name}'S", t, flags=re.IGNORECASE)
                        t = re.sub(r"1ST NAME", f"{first_name}'S", t, flags=re.IGNORECASE)
                        # Fix double possessive
                        t = t.replace("'S'S", "'S").replace("'S'S", "'S")
                        r.text = t.upper()
                        # Preserve font if it exists, otherwise set default
                        if not r.font.name:
                            r.font.name = "Schibsted Grotesk Medium"
                else:
                    # No runs, update paragraph text directly
                    t = p.text or ""
                    t = re.sub(r"1ST NAME[’']?S", f"{first_name}'S", t, flags=re.IGNORECASE)
                    t = re.sub(r"1ST NAME", f"{first_name}'S", t, flags=re.IGNORECASE)
                    t = t.replace("'S'S", "'S").replace("'S'S", "'S")
                    p.text = t.upper()
    
    # Step 6: Replace slide 11 (index 10) with PDF image if provided
    if image_bytes and len(prs.slides) > 10:
        add_or_replace_image(prs, 10, image_bytes)
    
    # Save to bytes
    out = BytesIO()
    prs.save(out)
    return out.getvalue()


# ============================================================================
# Google Slides Upload Functions
# ============================================================================

def cleanup_old_decks(drive, parent_folder_id: Optional[str] = None, max_files: int = 10):
    """Clean up old deck files to prevent quota issues."""
    try:
        query = "name contains 'Deck -' and mimeType='application/vnd.google-apps.presentation'"
        if parent_folder_id:
            query += f" and '{parent_folder_id}' in parents"
        
        list_kwargs = {
            "q": query,
            "orderBy": "createdTime desc",
            "fields": "files(id,name,createdTime)",
        }
        
        # Check if folder is in a Shared Drive
        shared_drive_id = None
        if parent_folder_id:
            try:
                folder_info = drive.files().get(
                    fileId=parent_folder_id,
                    fields="driveId",
                    supportsAllDrives=True
                ).execute()
                folder_drive_id = folder_info.get("driveId")
                if folder_drive_id:
                    shared_drive_id = folder_drive_id
                    list_kwargs.update({
                        "supportsAllDrives": True,
                        "includeItemsFromAllDrives": True,
                        "corpora": "drive",
                        "driveId": shared_drive_id,
                    })
            except Exception:
                pass
        
        results = drive.files().list(**list_kwargs).execute()
        files = results.get('files', [])
        
        if len(files) > max_files:
            files_to_delete = files[max_files:]
            for file in files_to_delete:
                try:
                    delete_kwargs = {"fileId": file['id']}
                    if shared_drive_id:
                        delete_kwargs["supportsAllDrives"] = True
                    drive.files().delete(**delete_kwargs).execute()
                    st.info(f"Cleaned up old deck: {file['name']}")
                except Exception as e:
                    st.warning(f"Could not delete {file['name']}: {e}")
    except Exception as e:
        st.warning(f"Could not clean up old files: {e}")


def upload_deck_to_google_drive(
    deck_type: str,
    client_name: str,
    priority_slide_index: Optional[int],
    service_column_slide_index: Optional[int],
    image_bytes: bytes | None = None,
    gallery_images: Optional[List[Dict[str, bytes]]] = None,
    pptx_bytes: bytes | None = None
) -> str:
    """Upload deck PPTX to Google Drive and return file URL. Similar to upload_workbook_to_google_sheet."""
    if not build or not MediaIoBaseUpload:
        raise RuntimeError("Google API client not available")
    
    # Get credentials
    creds = _get_credentials()
    drive = build("drive", "v3", credentials=creds)
    
    # Get parent folder ID for decks from secrets
    cfg = st.secrets.get("google_drive", {}) or {}
    parent_folder_id = cfg.get("parent_folder_id_deck")
    shared_drive_id = cfg.get("shared_drive_id")
    
    # Clean up old files
    cleanup_old_decks(drive, parent_folder_id)
    
    # Generate PPTX if not provided
    if pptx_bytes is None:
        pptx_bytes = create_deck_from_template(
            deck_type,
            client_name,
            priority_slide_index,
            service_column_slide_index,
            image_bytes,
            gallery_images
        )
    
    # Check if parent_folder_id is in a Shared Drive
    if parent_folder_id and not shared_drive_id:
        try:
            folder_info = drive.files().get(
                fileId=parent_folder_id,
                fields="id, name, driveId",
                supportsAllDrives=True
            ).execute()
            folder_drive_id = folder_info.get("driveId")
            if folder_drive_id:
                shared_drive_id = folder_drive_id
        except Exception:
            pass
    
    # Upload to Google Drive as PPTX file
    media = MediaIoBaseUpload(
        BytesIO(pptx_bytes),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        resumable=False
    )
    
    deck_name = f"Deck - {client_name or 'Client'} ({deck_type})"
    file_metadata = {
        "name": deck_name,
        "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    }
    
    if parent_folder_id:
        file_metadata["parents"] = [parent_folder_id]
    elif shared_drive_id:
        file_metadata["parents"] = [shared_drive_id]
    
    try:
        create_kwargs = {
            "body": file_metadata,
            "media_body": media,
            "fields": "id, webViewLink"
        }
        if shared_drive_id:
            create_kwargs["supportsAllDrives"] = True
        
        uploaded_file = drive.files().create(**create_kwargs).execute()
        file_id = uploaded_file.get("id")
        web_view = uploaded_file.get("webViewLink")
        
        if not file_id:
            raise RuntimeError("Upload succeeded but Google Drive did not return a file ID.")
        
        return web_view or f"https://drive.google.com/file/d/{file_id}/view"
    except Exception as e:
        error_msg = str(e)
        if "quota" in error_msg.lower():
            raise RuntimeError("Google Drive storage quota has been exceeded. Please delete older files or empty the Drive trash, then try again.")
        raise RuntimeError(f"Google Drive upload failed: {error_msg}")


def get_board_id_from_item(item_id: str) -> str | None:
    """Get board ID from a Monday.com item ID."""
    try:
        monday_config = st.secrets.get("monday", {})
        api_token = monday_config.get("api_token")
        
        if not api_token:
            return None
        
        query = f"""
        query {{
            items(ids: [{item_id}]) {{
                board {{
                    id
                }}
            }}
        }}
        """
        
        url = "https://api.monday.com/v2"
        headers = {
            "Authorization": api_token,
            "Content-Type": "application/json",
        }
        
        response = requests.post(url, json={"query": query}, headers=headers, timeout=30)
        data = response.json()
        
        if "errors" in data:
            return None
        
        items = data.get("data", {}).get("items", [])
        if not items:
            return None
        
        board = items[0].get("board", {})
        return board.get("id")
        
    except Exception:
        return None


def update_monday_item_deck_url(item_id: str, deck_url: str) -> bool:
    """Update a monday.com item with the deck URL in a 'Deck Link' field. Copied from update_monday_item_workbook_url."""
    try:
        monday_config = st.secrets.get("monday", {})
        api_token = monday_config.get("api_token")
        
        if not api_token:
            st.error("Monday.com API token not found in secrets.")
            return False
        
        # Query to get board columns to find the "Deck Link" column
        query = f"""
        query {{
            items(ids: [{item_id}]) {{
                board {{
                    id
                    columns {{
                        id
                        title
                        type
                    }}
                }}
            }}
        }}
        """
        
        url = "https://api.monday.com/v2"
        headers = {
            "Authorization": api_token,
            "Content-Type": "application/json",
        }
        
        response = requests.post(url, json={"query": query}, headers=headers, timeout=30)
        data = response.json()
        
        if "errors" in data:
            st.error(f"Error fetching monday.com columns: {data['errors']}")
            return False
        
        items = data.get("data", {}).get("items", [])
        if not items:
            st.error("Item not found in monday.com")
            return False
        
        board = items[0].get("board", {})
        board_id = board.get("id")
        columns = board.get("columns", [])
        
        if not board_id:
            st.error("Could not determine board ID from monday.com item")
            return False
        
        # Find "Deck Link" column
        deck_column = None
        for col in columns:
            title_lower = col.get("title", "").lower()
            if "deck" in title_lower and "link" in title_lower:
                deck_column = col
                break
        
        if not deck_column:
            st.warning("⚠️ 'Deck Link' column not found in monday.com. Please create a URL column named 'Deck Link' in the Sales board.")
            return False
        
        column_id = deck_column.get("id")
        column_type = deck_column.get("type")
        
        # Update the item with the deck URL
        # For URL columns, the value format is: {"url": "https://...", "text": "Link Text"}
        if column_type == "link":
            mutation = f"""
            mutation {{
                change_column_value(
                    board_id: {board_id},
                    item_id: {item_id},
                    column_id: "{column_id}",
                    value: "{{\\"url\\": \\"{deck_url}\\", \\"text\\": \\"View Deck\\"}}"
                ) {{
                    id
                }}
            }}
            """
        else:
            # For text columns, just use the URL as text
            mutation = f"""
            mutation {{
                change_column_value(
                    board_id: {board_id},
                    item_id: {item_id},
                    column_id: "{column_id}",
                    value: "{deck_url}"
                ) {{
                    id
                }}
            }}
            """
        
        response = requests.post(url, json={"query": mutation}, headers=headers, timeout=30)
        result = response.json()
        
        if "errors" in result:
            st.error(f"Error updating monday.com: {result['errors']}")
            return False
        
        return True
        
    except Exception as e:
        st.error(f"Failed to update monday.com: {e}")
        return False


# ============================================================================
# UI Components
# ============================================================================

def render_deck_type_selector() -> str:
    """Render deck type selector and return selected type."""
    st.subheader("Deck Type")
    deck_type = st.selectbox(
        "Select Deck Type",
        options=list(DECK_TYPES.keys()),
        index=0,
        key="dc_deck_type",
        help="Choose the deck template style"
    )
    return deck_type


def render_proposal_for_input(default_value: str = "") -> str:
    """Render proposal for input and return client name.
    
    Args:
        default_value: Default value for the client name input (from query params)
    """
    st.subheader("Proposal For")
    # Use default_value if provided, otherwise use session state
    initial_value = default_value if default_value else st.session_state.get("dc_client_name", "")
    client_name = st.text_input(
        "Client Name",
        value=initial_value,
        key="dc_client_name",
        help="Enter the client name for the 'Proposal For' slide",
        placeholder="Enter client name"
    )
    return client_name


def render_priorities_selector() -> Optional[int]:
    """Render priorities selector and return selected slide index."""
    st.subheader("Priorities")
    
    priorities, _ = extract_slides_from_service_columns()
    
    if not priorities:
        st.info("No priority slides found in SERVICE COLUMNS.pptx")
        return None
    
    # Create options for radio buttons using "Priorities 1", "Priorities 2", etc.
    options = [priorities[idx]["title"] for idx in sorted(priorities.keys())]
    selected = st.radio(
        "Select Priorities",
        options=options,
        index=0,
        key="dc_priority_selection",
        help="Choose one of the pre-formatted Priority slides (all or nothing selection)"
    )
    
    # Show expanders with sub-fields for each priority (similar to Additional Packages)
    for idx in sorted(priorities.keys()):
        priority_data = priorities[idx]
        with st.expander(f"{priority_data['title']} - sub-fields", expanded=False):
            if priority_data.get("sub_fields"):
                sub_fields_text = "\n".join([f"- {field}" for field in priority_data["sub_fields"]])
                st.markdown(sub_fields_text)
            else:
                st.info("No sub-fields found for this priority")
    
    # Extract slide index from selection
    for idx, priority_data in priorities.items():
        if priority_data["title"] == selected:
            return idx
    
    # Default to first priority if somehow no match
    return list(priorities.keys())[0] if priorities else None


def render_service_columns_selector(deck_type: str) -> Optional[int]:
    """Render service columns selector and return selected slide index.
    
    Args:
        deck_type: Selected deck type
    """
    st.subheader("Services")
    
    _, service_columns = extract_slides_from_service_columns()
    
    if not service_columns:
        st.info("No service column slides found in SERVICE COLUMNS.pptx")
        return None
    
    # Define the 9 specific service column options (slides 0-8, which are indices 0-8)
    # These correspond to slides 1-9 in the PowerPoint (0-based indexing)
    allowed_indices = [0, 1, 2, 3, 4, 5, 6, 7, 8]
    
    # Define the formatted titles with (Style) and (Service) labels for each slide
    formatted_titles = {
        0: "SOURCING(Style)-DEVELOPMENT",
        1: "SOURCING(Style)-DESIGN(Service)-DEVELOPMENT(Style)",
        2: "SOURCING(Style)-WASH/DYE(Style)-DEVELOPMENT(Style)",
        3: "DESIGN(Style)-TREATMENT(Service)-SOURCING(Style)-DEVELOPMENT",
        4: "SOURCING(Style)-DESIGN(Style)-DEVELOPMENT(Style)",
        5: "SOURCING(Style)-WASH/DYE(Style)",
        6: "SOURCING(Style)-FABRIC TREATMENT(Style)",
        7: "SOURCING(Style)-TREATMENT(Service)-DEVELOPMENT",
        8: "TREATMENT(Service)-SOURCING(Style)-DEVELOPMENT"
    }
    
    # Build options from allowed indices only, using formatted titles
    options_map = {}  # formatted_title -> index
    options_list = []
    
    for idx in allowed_indices:
        if idx in service_columns:
            # Use the formatted title for display
            formatted_title = formatted_titles.get(idx, service_columns[idx]["title"])
            options_map[formatted_title] = idx
            options_list.append(formatted_title)
    
    if not options_list:
        st.info("No valid service column slides found. Please ensure SERVICE COLUMNS.pptx has slides 1-9.")
        return None
    
    selected = st.selectbox(
        "Select Services",
        options=options_list,
        index=0,
        key="dc_service_column_selection",
        help="Choose one of the pre-formatted Service Column slides"
    )
    
    # Return the slide index for the selected title
    return options_map.get(selected)


def render_pdf_upload() -> bytes | None:
    """Render PDF upload and return image bytes from first page."""
    st.subheader("PDF Upload")
    pdf_file = st.file_uploader("Upload a PDF to insert as first-page image", type=["pdf"], key="dc_pdf_upload")
    img_bytes = st.session_state.get('dc_img_bytes')
    if pdf_file is not None:
        try:
            import fitz  # PyMuPDF
            pdf_bytes = pdf_file.read()
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            if doc.page_count > 0:
                page = doc.load_page(0)
                pix = page.get_pixmap(dpi=180)
                img_bytes = pix.tobytes("png")
                st.session_state['dc_img_bytes'] = img_bytes
                st.success("PDF converted to image")
            doc.close()
        except Exception as e:
            st.warning(f"Could not convert PDF to image: {e}")
            img_bytes = None
    elif img_bytes is not None:
        st.info("Using previously uploaded PDF image (restored)")
    
    return img_bytes


def _extract_clean_filename(display_name: str) -> str:
    """Extract clean filename without path and extension.
    
    Examples:
        "images/ALABAMA BLONDE 1.webp" -> "ALABAMA BLONDE 1"
        "folder/subfolder/image.png" -> "image"
    """
    # Get just the filename (last part after /)
    filename = display_name.split("/")[-1]
    # Remove extension
    if "." in filename:
        filename = filename.rsplit(".", 1)[0]
    return filename


def _remove_trailing_numbers(name: str) -> str:
    """Remove trailing numbers from a name.
    
    Examples:
        "ALABAMA BLONDE 4" -> "ALABAMA BLONDE"
        "CAROLINE HAYDEN 1" -> "CAROLINE HAYDEN"
        "IMAGE 123" -> "IMAGE"
    """
    import re
    # Remove trailing spaces and numbers
    return re.sub(r'\s+\d+$', '', name.strip())


def render_gallery_selector() -> List[Dict[str, Any]]:
    """Render the gallery selector (exactly 3 images) below the PDF upload."""
    
    # Add visual preview option - hyperlink to Google Drive folder
    GALLERY_DRIVE_URL = f"https://drive.google.com/drive/folders/{IMAGE_GALLERY_FOLDER_ID}?usp=drive_link"
    
    # Header with refresh button
    col1, col2 = st.columns([1, 1])
    with col1:
        st.subheader(f"[Gallery Images]({GALLERY_DRIVE_URL})")
    with col2:
        if st.button("🔄 Refresh Gallery", help="Clear cache and reload images from Google Drive"):
            # Clear all cached data to force reload from Google Drive
            st.cache_data.clear()
            # Also clear session state cache if any
            if "gallery_images_cache" in st.session_state:
                del st.session_state["gallery_images_cache"]
            st.success("Gallery refreshed! Reloading...")
            st.rerun()

    if not Image:
        st.error("Pillow is required to load gallery images. Please install pillow.")
        return []
    
    try:
        gallery_items = list_gallery_images()
    except Exception as e:
        st.warning(f"Could not load gallery images: {e}")
        return []
    
    if not gallery_items:
        st.info("No images found in the gallery folder.")
        return []
    
    # Create a mapping of clean filenames to items for easy lookup
    clean_name_to_item = {}
    for item in gallery_items:
        display_name = item.get("display_name") or item.get("name") or "Image"
        clean_name = _extract_clean_filename(display_name)
        clean_name_to_item[clean_name] = item
    
    # Get list of all clean image names for the multiselect
    image_names = sorted(clean_name_to_item.keys())
    
    # Single multiselect input - let Streamlit manage the widget state completely
    widget_key = "dc_gallery_multiselect"
    selected_names = st.multiselect(
        "Select 3 Images",
        options=image_names,
        key=widget_key,
        help="Type to search and select exactly 3 images from the gallery"
    )
    
    # Filter out any invalid selections (shouldn't happen, but safety check)
    selected_names = [name for name in selected_names if name in image_names]
    
    # Limit to 3 selections (in case user somehow selects more)
    if len(selected_names) > 3:
        selected_names = selected_names[:3]
        st.warning("⚠️ Only 3 images can be selected. Showing the first 3.")
    
    # Track selection for caching - use a simple comparison
    prev_selection_for_cache = st.session_state.get("dc_prev_gallery_selection", None)
    selection_changed = prev_selection_for_cache != selected_names
    if selection_changed:
        st.session_state["dc_prev_gallery_selection"] = selected_names.copy() if selected_names else []
    
    # Only show images when exactly 3 are selected
    if len(selected_names) != 3:
        if len(selected_names) > 0:
            st.info(f"Please select exactly 3 images ({len(selected_names)}/3 selected).")
        else:
            st.info("Please select exactly 3 images to continue.")
        # Clear cache if selection is incomplete
        if "dc_cached_gallery_images" in st.session_state:
            del st.session_state["dc_cached_gallery_images"]
        if "dc_cached_gallery_selection" in st.session_state:
            del st.session_state["dc_cached_gallery_selection"]
        return []
    
    # Check if we have cached images for this exact selection
    cache_key = tuple(selected_names)  # Use tuple as cache key (order matters)
    cached_images_key = "dc_cached_gallery_images"
    cached_selection_key = "dc_cached_gallery_selection"
    
    # Load images only if selection changed or cache is invalid
    if (selection_changed or 
        st.session_state.get(cached_selection_key) != cache_key or
        cached_images_key not in st.session_state):
        
        # Load and cache the images
        selected_images: List[Dict[str, Any]] = []
        for clean_name in selected_names:
            if clean_name not in clean_name_to_item:
                st.error(f"❌ Image '{clean_name}' not found in gallery.")
                return []
            
            item = clean_name_to_item[clean_name]
            try:
                img_bytes = get_resized_gallery_image(item["id"])
                selected_images.append({
                    "id": item["id"],
                    "name": clean_name,
                    "data": img_bytes
                })
            except Exception as e:
                st.error(f"❌ Could not load image '{clean_name}': {e}")
                return []
        
        # Cache the images
        st.session_state[cached_images_key] = selected_images
        st.session_state[cached_selection_key] = cache_key
    else:
        # Use cached images
        selected_images = st.session_state[cached_images_key]
    
    # Display the images
    cols = st.columns(3)
    for idx, img_data in enumerate(selected_images):
        with cols[idx]:
            st.image(img_data["data"], use_container_width=True)
            # Remove trailing numbers from displayed name
            display_name = _remove_trailing_numbers(img_data["name"])
            st.caption(display_name)
    
    return selected_images


# ============================================================================
# Main Application
# ============================================================================

def main():
    st.title("📽️ Deck Creator")
    
    # Get query parameters from Monday.com link
    query_params = st.query_params
    first_name = query_params.get("first_name", "").strip()
    last_name = query_params.get("last_name", "").strip()
    item_id = query_params.get("item_id", "").strip()
    
    # Combine first_name and last_name for client_name default
    client_name_default = ""
    if first_name or last_name:
        client_name_default = f"{first_name} {last_name}".strip()
    
    # Render UI components
    deck_type = render_deck_type_selector()
    client_name = render_proposal_for_input(client_name_default)
    priority_slide_index = render_priorities_selector()
    service_column_slide_index = render_service_columns_selector(deck_type)
    image_bytes = render_pdf_upload()
    gallery_images = render_gallery_selector()
    
    # Show Monday.com upload section
    st.markdown("---")
    st.subheader("Google Sheets --> Monday.com Upload")
    st.caption("Uploads will use the shared Google Drive folder configured for the service account.")
    
    # Action buttons
    st.divider()
    col_dl, col_gs = st.columns(2)
    
    with col_dl:
        generate_button = st.button("Create PowerPoint", type="primary")
    
    with col_gs:
        push_gslides_button = st.button("Google Sheets --> Monday.com Upload", type="primary")
    
    # Handle PowerPoint generation
    if generate_button:
        if not Presentation:
            st.error("python-pptx is not installed.")
            return
        if len(gallery_images) != 3:
            st.error("Please select exactly 3 gallery images before generating.")
            return
        
        try:
            pptx_bytes = create_deck_from_template(
                deck_type,
                client_name,
                priority_slide_index,
                service_column_slide_index,
                image_bytes,
                gallery_images
            )
            
            # Store in session state and trigger auto-download
            filename = f"deck_{client_name or 'client'}_{deck_type.replace(' ', '_')}.pptx"
            st.session_state['pptx_download_data'] = pptx_bytes
            st.session_state['pptx_download_filename'] = filename
            st.session_state['pptx_download_timestamp'] = datetime.now().timestamp()
            st.session_state['pptx_auto_download'] = True
            st.success("PowerPoint generated and downloaded!!!")
            
        except Exception as e:
            st.error(f"Failed to generate PowerPoint: {e}")
            import traceback
            st.code(traceback.format_exc())
    
    # Auto-download logic
    if st.session_state.get('pptx_auto_download', False):
        pptx_bytes = st.session_state.get('pptx_download_data')
        filename = st.session_state.get('pptx_download_filename', 'deck.pptx')
        
        if pptx_bytes:
            st.session_state['pptx_auto_download'] = False
            
            b64_pptx = base64.b64encode(pptx_bytes).decode()
            timestamp = st.session_state.get('pptx_download_timestamp', 'default')
            
            import streamlit.components.v1 as components
            
            components.html(f"""
            <script>
            (function() {{
                const downloadId = 'pptx_download_{timestamp}';
                if (window.pptxDownloadTriggered === downloadId) {{
                    return;
                }}
                window.pptxDownloadTriggered = downloadId;
                
                function triggerDownload() {{
                    try {{
                        const data = 'data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_pptx}';
                        const link = document.createElement('a');
                        link.href = data;
                        link.download = '{filename}';
                        link.style.display = 'none';
                        document.body.appendChild(link);
                        link.click();
                        setTimeout(() => {{
                            if (link.parentNode) {{
                                document.body.removeChild(link);
                            }}
                        }}, 100);
                        return true;
                    }} catch(e) {{
                        console.error('Download failed:', e);
                        return false;
                    }}
                }}
                
                setTimeout(triggerDownload, 100);
            }})();
            </script>
            """, height=0)
            
            st.success(f"✅ PowerPoint generated: {filename}")
    
    # Handle Google Drive --> Monday.com upload
    if push_gslides_button:
        if not build:
            st.error("Google API client not available")
        else:
            if len(gallery_images) != 3:
                st.error("Please select exactly 3 gallery images before uploading.")
                return
            with st.spinner("Uploading deck to Google Drive and updating Monday.com..."):
                try:
                    # Generate PPTX bytes first (needed for both Google Drive and Monday.com upload)
                    pptx_bytes = create_deck_from_template(
                        deck_type,
                        client_name,
                        priority_slide_index,
                        service_column_slide_index,
                        image_bytes,
                        gallery_images
                    )
                    
                    # Upload to Google Drive (pass PPTX bytes to avoid regenerating)
                    deck_url = upload_deck_to_google_drive(
                        deck_type,
                        client_name,
                        priority_slide_index,
                        service_column_slide_index,
                        image_bytes,
                        gallery_images,
                        pptx_bytes
                    )
                    st.success(f"✅ Deck uploaded to Google Drive: [Open Deck]({deck_url})")
                    
                    # Update Monday.com with the deck URL if item_id is provided
                    if item_id:
                        # Update Deck Link column
                        deck_link_success = update_monday_item_deck_url(item_id, deck_url)
                        
                        # Show results
                        if deck_link_success:
                            st.success(f"✅ Monday.com item updated with Deck Link!")
                        else:
                            st.warning("⚠️ Deck uploaded to Google Drive, but failed to update Monday.com item. Please update manually.")
                    else:
                        st.info("ℹ️ No Monday.com item ID provided. Deck uploaded to Google Drive only.")
                    
                except Exception as e:
                    error_msg = str(e)
                    if "quota" in error_msg.lower():
                        st.error("❌ Google Drive quota exceeded. Please try again later or check your Google Drive storage.")
                        st.markdown("""
                        **To resolve this issue:**
                        - Free up space in your Google Drive
                        - Delete old deck files from Google Drive
                        - Try using a different Google account with more storage
                        - Wait a few hours and try again (quotas reset periodically)
                        """)
                    elif "permission" in error_msg.lower():
                        st.error("❌ Permission denied. Please check your Google API credentials and permissions.")
                    elif "not found" in error_msg.lower():
                        st.error("❌ Template file not found. Please ensure the template PPTX file exists in the inputs folder.")
                    else:
                        st.error(f"❌ Failed to upload deck: {error_msg}")


if __name__ == "__main__":
    main()